import pandas as pd
from datetime import datetime, date
import matplotlib.pyplot as plt
import collections
from collections import Counter
import calendar
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import os
import sys
import logging

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

yyyy,mm=2017,12
logger.info('Starting run for %i-%i...' % (yyyy,mm))
print("Hello world")


# Colour codes for HPE accent and default chart colours
TURQUOISE = '#2AD2C9'
PURPLE    = '#5C4767'
ORANGE    = '#FF8D6D'
DARK_STEEL= '#5F7A6C'
GRAY      = '#C6C9CA'
DARK_GRAY = '#808285'

def make_date(cell_val):
    if type(cell_val) is datetime:
        v=cell_val.date()
    elif type(cell_val) is str: 
        try:
            v= datetime.strptime(cell_val,"%b %d, %Y").date()
        except ValueError:
            v=date.fromordinal(1)
    else:
        v=date.fromordinal(2)
    return v

def replace_strings(series,repl_dict):
    for k,v in repl_dict.items():
        series.str.replace(k,v)
    return(series)

def count_strings(list_of_strings,s_to_find):
    """Count how many times the strings in sequence s occur in the strings in list_of_strings 
       Return a Counter with the number of times each i in s occurs in one of the strings in list_of_strings
    """
    if s_to_find is str:
        count_list=sum(list_of_strings.str.find(str)>0)
    else:
        count_list=Counter()
        for i in s_to_find:
            count_list[i]=sum(list_of_strings.str.find(i)>0)
    return count_list
        
def sum_of_dicts(d1,d2):
    dd_sum={}
    for k in set(d1.keys()).union(d2.keys()):
        dd_sum[k]=0
        if k in d1:
            dd_sum[k]+=d1[k]
        if k in d2:
            dd_sum[k]+=d2[k]
        if (dd_sum[k]==0): del dd_sum[k]
    return(dd_sum)

def keyword_counts_in_series(series,keywords):
    series = replace_strings(series.str.lower().str.replace('[\n&,.-]',' '),
                             replace_text
                            )
    return(count_strings(series,keywords))

def dataframe_for_month(df, year=2017, month=1):
    """Yields a subset the dataframe with only those rows in the given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    if (mm==12): mm2,yyyy2=1,yyyy+1
    else: mm2,yyyy2=mm+1,yyyy
    
    logger.debug("Selecting rows for %i-%i" % (yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy,mm,1)) & (df.date<date(yyyy2,mm2,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))
    
    return month_df


def keywords_in_dataframe(month_df,year=2017,month=1):
    """ Count each keyword in the dataframe's important columns for a given year & month
        Returns a Counter collection of (keyword:count)
    """
    import operator   #for sorted(), used in sorting items into OrderedDict
     
    logger.debug("Counting occurrences of keywords for this month")
    c_wtlma = keyword_counts_in_series(month_df['Want to Learn More About'],vocab)
    c_actions=keyword_counts_in_series(month_df['Action Items'],vocab)

    logger.debug("Adding together the results")
    counter_words_to_freq = c_wtlma + c_actions
    
    return counter_words_to_freq

def count_rows_with_comments(df):
    """ Count the number of rows in dataframe with a comment in either <Want to Learn More About> or <Action Items>
        Returns the count
    """
    return (df["Want to Learn More About"].notnull() | df["Action Items"].notnull()).sum()

def counts_by_centre(dataframe):
    centre_counts = dataframe.Ctr.value_counts()
    counts_list = (centre_counts.PA if hasattr(centre_counts,"PA") else 0, 
                   centre_counts.H if hasattr(centre_counts,"H") else 0, 
                   centre_counts.NY1 if hasattr(centre_counts,"NY1") else 0, 
                   centre_counts.LON1 if hasattr(centre_counts,"LON1") else 0, 
                   centre_counts.SNG if hasattr(centre_counts,"SNG") else 0
                   )
    return counts_list

def file_wordcloud_for_month(keywords_for_month, year=2017, month=8):
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    # (to install, run "pip install wordcloud"
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt
    import os
    
    font_path = os.path.join("C:",os.sep,"Windows",os.sep,"Fonts",os.sep,'arial.ttf')
    logger.debug("Generating the wordcloud")
    wc_for_month = WordCloud(font_path=font_path,
                             width=2500,height=500,
                             prefer_horizontal=1.0,
                             relative_scaling=0.7,
                             max_font_size=196,
                             background_color="white"
                             ).generate_from_frequencies(keywords_for_month)
    
    # Display the generated image
    plt.imshow(wc_for_month,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    #plt.show()
    
    filename = "wordcloud-"+calendar.month_name[month]+".png"
    plt.imsave(filename,wc_for_month,format="png")
    plt.close()
    
    return filename


def file_graph_for_month_kwd(kwd,kwd_pos,vals,months,line_color):
    fig,ax=plt.subplots()
    ax.clear()
    ax.plot(vals,line_color,linewidth=3)
    #ax.set_title(kwd,fontsize=48)
    ax.set_axis_off()
    m_minus_2_percent = "{0:.0f}%".format(vals[0] * 100)
    m_this_percent    = "{0:.0f}%".format(vals[2] * 100)
    ax.text(0,vals[0],calendar.month_name[months[0]]+"\n"+m_minus_2_percent,fontsize=36)
    ax.text(2,vals[2],calendar.month_name[months[2]]+"\n"+m_this_percent,fontsize=36)
    filename = "graph-"+str(kwd_pos)+".png"
    logger.debug("Saving %s graph for keyword %s in file %s" % (kwd_pos,kwd,filename))
    fig.savefig(filename,bbox_inches="tight")
    plt.close(fig)
    
    return filename

def file_donut_pie_for_month(values,kwd):
    fig, ax = plt.subplots()
    ax.axis('equal')
    width = 0.40
    outside, _ = ax.pie(values,startangle=90,counterclock=False,colors=(TURQUOISE, PURPLE, ORANGE, DARK_STEEL, GRAY))
    ax.legend(("Palo Alto","Houston","NY","London","Singapore"),fontsize=36,bbox_to_anchor=(0.8,1.0),frameon=False)
    plt.setp( outside, width=width, edgecolor='white')
    
    filename = "donut-"+str(kwd)+".png"
    logger.debug("Saving %s donut in file %s with values %r" % (kwd, filename, values))
    fig.savefig(filename,bbox_inches="tight")
    
    plt.close(fig)
    
    return filename

def file_donut_pie_for_industries(industries):
    fig, ax = plt.subplots()
    ax.axis('equal')
    width = 2.0
    outside, _ = ax.pie(industries.values,
                        startangle=90,counterclock=False,
                        colors=(TURQUOISE, PURPLE, ORANGE, DARK_STEEL, GRAY,DARK_GRAY),
                        radius=5.0)
    ax.legend(industries.index.tolist(),ncol=3,fontsize=24,loc=10,bbox_to_anchor=(0.5,-2),frameon=False)
    plt.setp( outside, width=width, edgecolor='white')
    
    filename = "donut-industries.png"
    logger.debug("Saving Industries donut in file %s with values %r" % (filename, industries))
    fig.savefig(filename,bbox_inches="tight")
    
    plt.close(fig)
    
    return filename

def write_customer_list(df_for_kwd,text_frame):    
    text_frame.text = "Customers"
    font=text_frame.paragraphs[0].font
    font.size = Pt(12)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1
    
    for i in (df_for_kwd)["Account Name"].str.title().iteritems(): 
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        new_para.text=i[1]
    
    return  

def write_top_keywords(text_frame,titleText,kwd_counts,rowcount):
    text_frame.text = titleText
    font=text_frame.paragraphs[0].font
    font.size = Pt(10)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1
    
    n=0
    for k,v in kwd_counts.items():
        kwd_counts[k]=v/rowcount
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        new_para.text=k+" - "+"{0:.0f}%".format(v/rowcount * 100)
        n+=1
        if (n>=4): break;
    
    return

def find_text_in_shapes(slide_shapes,searchword):
    """Find the given text in the list of powerpoint shapes passed in slide_shapes
    """
    found_idx = -1
    for i in range(len(slide_shapes)):
        if (slide_shapes[i].has_text_frame):
            text_frame = slide_shapes[i].text_frame 
            if (text_frame.text.find(searchword)>=0):
                found_idx = i
                break
            
    return found_idx

def new_run_in_slide(para,text='text',fontname='Arial',fontsize=24):
    "Add new run to a paragraph in a powerpoint slide"
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.name = fontname
    return run

replace_text={
    "/": " ",                       "\n": " ",                             "big data": "big_data",    
    "new stack": "new_stack",       "flexible capacity": "flex_capacity",  "flex capacity": "flex_capacity",
    "gen z": "gen-z",               "cloud cruiser": "cloudcruiser",       "store once": "storeonce",
    "smart city": "smart_city",     "future city": "smart_city",           "azure stack": "azurestack",
    "intelligent edge": "edgeline", "clear pass": "clearpass",             "one view": "oneview",
    "the machine": "the_machine",   "open stack": "openstack",             "store once": "storeonce",
    "office 365": "office365",      "hp financial services":"hpefs",       "mobility": "wireless",
    "hpe financial services": "hpefs",     "hpfs": "hpefs",                "integrity": "superdome",
    "one sphere": "onesphere",      "gen 10": "gen10"
} 

#use this specific set of words as the vocab to look for
vocab=["iot","edgeline","smart_city", "big_data","sap","apollo", "saas", "analytics", "sgi", 
       "networking", "wireless", "aruba", "arista", "clearpass", "naas", "skype", "meridian", 
       "airwave","composable", "niara", "flex_capacity", "helion", "azurestack","synergy",
       "moonshot","converged","hyperconverged","hybrid","docker","vdi","new_stack","cloudline",
       "cloudsystem","easyconnect","openstack","devops","3par","bura","simplivity","nimble",
       "scality","storeonce", "oneview", "office365","pointnext","hpefs","the_machine","photonics",
       "blockchain","pathfinder","gen10","gen-z","cloudcruiser","blades","superdome",
       "onesphere","greenlake"
      ]

excel_file='Insights.xlsx'
logger.info("Importing excel file "+excel_file)
all_df = pd.read_excel(open(excel_file,'rb'),header=8)

logger.debug("Adding structured date column")
all_df.insert(loc=0,column='date',value=all_df['Visit Date'].apply(make_date))

## Given the month and year, calc the number of the previous two months
if (mm==1): mm_minus_1,year_for_mm_minus_1 = 12, yyyy-1
else: mm_minus_1, year_for_mm_minus_1 = mm-1, yyyy
if (mm<=2): mm_minus_2,year_for_mm_minus_2 = mm+10, yyyy-1
else: mm_minus_2, year_for_mm_minus_2 = mm-2, yyyy   
    
### Now start to generate the powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt

# These appear to be the layouts used in the master for this slide deck
LAYOUT_TITLE_WITH_DARK_PICTURE = 0
LAYOUT_TITLE_SLIDE_WITH_NAME   = 1
LAYOUT_DIVIDER                 = 2
LAYOUT_TITLE_ONLY              = 3
LAYOUT_TITLE_AND_SUBTITLE      = 4
LAYOUT_BLANK                   = 5

## Open up the source presentation    
prs = Presentation('GCA_Customer_Insights_Month-Year.pptx')
this_month = calendar.month_name[mm]
earliest_month = calendar.month_name[mm_minus_2]
logger.info("Creating presentation for %s" % (this_month))

################################################
## Modify existing title slide with month
################################################
logger.info("Modifying text in slide 0")
s = prs.slides[0]
slide_shapes = s.shapes
text_frame = slide_shapes[0].text_frame  # should be the the title textframe
#clear existing text, and write new text into title textframe
text_frame.clear() 
# First para with a run of 60-point Arial font
new_run_in_slide(text_frame.paragraphs[0],text='Customer Insights',fontname='Arial',fontsize=60)  
# Second para with a run of 32-point font
new_run_in_slide(text_frame.add_paragraph(),text='Learnings from '+this_month+' EBC/CEC visits',fontsize=32)    

################################################
## 3rd slide: count the keywords for this month and build the wordcloud
################################################
df_for_month = dataframe_for_month(all_df, year=yyyy, month=mm)
kwd_count_for_month = keywords_in_dataframe(df_for_month,year=yyyy,month=mm)
logger.info("Top keywords and counts found for month %i : %r" % (mm,kwd_count_for_month.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_month,
                         year=yyyy,month=mm)
## Modifying main wordcloud slide by changing title and adding pic for this month's wordcloud
logger.info("Modifying text and adding wordcloud in slide 2")
s = prs.slides[2]
slide_shapes=s.shapes
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="In "+this_month+" customers wanted to learn more about...",
       fontname="Arial",fontsize=28)
left=Inches(0.5); top=Inches(4.0)
slide_shapes.add_picture("wordcloud-"+this_month+".png",left,top,height=Inches(2.5))

################################################
## 4th slide: count the keywords for previous two months, and build their wordclouds.
################################################
df_for_month_minus_1 = dataframe_for_month(all_df, 
                                           year=year_for_mm_minus_1, month=mm_minus_1)
kwd_count_for_m_minus_1 = keywords_in_dataframe(df_for_month_minus_1,
                                                   year=year_for_mm_minus_1,month=mm_minus_1)
logger.info("Top keywords and counts found for month %i : %r" % (mm_minus_1,kwd_count_for_m_minus_1.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_m_minus_1,
                         year=year_for_mm_minus_1,month=mm_minus_1)

df_for_month_minus_2 = dataframe_for_month(all_df, 
                                           year=year_for_mm_minus_2, month=mm_minus_2)
kwd_count_for_m_minus_2 = keywords_in_dataframe(df_for_month_minus_2,
                                                   year=year_for_mm_minus_2,month=mm_minus_2)
logger.info("Top keywords and counts found for month %i : %r" % (mm_minus_2,kwd_count_for_m_minus_2.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_m_minus_2,
                         year=year_for_mm_minus_2,month=mm_minus_2)


## Modifying 3-month wordcloud slide by adding pic for last 3 months' wordcloud
logger.info("Adding three wordclouds in slide 3")
s = prs.slides[3]
slide_shapes=s.shapes
left=Inches(3.65)
top=Inches(1.4)
slide_shapes.add_picture("wordcloud-"+calendar.month_name[mm_minus_2]+".png",left,top,height=Inches(1.5))
top=Inches(3.25)
slide_shapes.add_picture("wordcloud-"+calendar.month_name[mm_minus_1]+".png",
                         left,top,height=Inches(1.5))
top=Inches(5.05)
slide_shapes.add_picture("wordcloud-"+this_month+".png",
                         left,top,height=Inches(1.5))
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
i = find_text_in_shapes(slide_shapes,"Month-2")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = calendar.month_name[mm_minus_2]
else:
    logger.error("Could not find Month-2 placeholder on 4th slide")
i = find_text_in_shapes(slide_shapes,"Month-1")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = calendar.month_name[mm_minus_1]
else:
    logger.error("Could not find Month-1 placeholder on 4th slide")
i = find_text_in_shapes(slide_shapes,"Month-0")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = calendar.month_name[mm]
else:
    logger.error("Could not find Month-0 placeholder on 4th slide")

################################################
## Get top 3 keywords for this month, and graph their usage 
################################################
top_3 = kwd_count_for_month.most_common(3)   # top 3 keywords for most recent month in list with their counts
kwd0=top_3[0][0]
kwd1=top_3[1][0]
kwd2=top_3[2][0]
logger.debug("Top 3 keywords for current month are %s %s %s" % (kwd0,kwd1,kwd2))
months=[mm_minus_2,mm_minus_1,mm]
useful_rows_in_m_2 = count_rows_with_comments(df_for_month_minus_2)
useful_rows_in_m_1 = count_rows_with_comments(df_for_month_minus_1)
useful_rows_in_m = count_rows_with_comments(df_for_month)

kwd0_c2 = kwd_count_for_m_minus_2[kwd0] if (kwd0 in kwd_count_for_m_minus_2) else 0
kwd0_c1 = kwd_count_for_m_minus_1[kwd0] if (kwd0 in kwd_count_for_m_minus_1) else 0
kwd0_c0 = kwd_count_for_month[kwd0]     # must have keyword as it came from this dictionary
vals_kwd0=[kwd0_c2/useful_rows_in_m_2, kwd0_c2/useful_rows_in_m_1, kwd0_c0/useful_rows_in_m]
file_linegraph_topic1 = file_graph_for_month_kwd(kwd0,"1st",vals_kwd0,months,TURQUOISE)

kwd1_c2 = kwd_count_for_m_minus_2[kwd1] if (kwd1 in kwd_count_for_m_minus_2) else 0
kwd1_c1 = kwd_count_for_m_minus_1[kwd1] if (kwd1 in kwd_count_for_m_minus_1) else 0
kwd1_c0 = kwd_count_for_month[kwd1]     # must have keyword as it came from this dictionary
vals_kwd1=[kwd1_c2/useful_rows_in_m_2, kwd1_c2/useful_rows_in_m_1, kwd1_c0/useful_rows_in_m]
file_linegraph_topic2 = file_graph_for_month_kwd(kwd1,"2nd",vals_kwd1,months,PURPLE)

kwd2_c2 = kwd_count_for_m_minus_2[kwd2] if (kwd2 in kwd_count_for_m_minus_2) else 0
kwd2_c1 = kwd_count_for_m_minus_1[kwd2] if (kwd2 in kwd_count_for_m_minus_1) else 0
kwd2_c0 = kwd_count_for_month[kwd2]     # must have keyword as it came from this dictionary
vals_kwd2=[kwd2_c2/useful_rows_in_m_2, kwd2_c2/useful_rows_in_m_1, kwd2_c0/useful_rows_in_m]
file_linegraph_topic3 = file_graph_for_month_kwd(kwd2,"3rd",vals_kwd2,months,ORANGE)

## Build a subset of the dataframe for last 3 months that uses each of the top 3 kwds in this month
df_for_3months = pd.concat([df_for_month,df_for_month_minus_1,df_for_month_minus_2])
df_for_kwd0 = df_for_3months.loc[(df_for_3months['Want to Learn More About'].str.lower().str.find(kwd0)>0) |
                                 (df_for_3months['Action Items'].str.lower().str.find(kwd0)>0) 
                                ]
df_for_kwd1 = df_for_3months.loc[(df_for_3months['Want to Learn More About'].str.lower().str.find(kwd1)>0) |
                                 (df_for_3months['Action Items'].str.lower().str.find(kwd1)>0) 
                                ]
df_for_kwd2 = df_for_3months.loc[(df_for_3months['Want to Learn More About'].str.lower().str.find(kwd2)>0) |
                                 (df_for_3months['Action Items'].str.lower().str.find(kwd2)>0) 
                                ]

## Create donut pies showing split of visits expressing interest in top 3 topics by centre over last 3 months
kwd0_counts = counts_by_centre(df_for_kwd0)
file_donut_topic1 = file_donut_pie_for_month(counts_by_centre(df_for_kwd0),"1st")
file_donut_topic2 = file_donut_pie_for_month(counts_by_centre(df_for_kwd1),"2nd")
file_donut_topic3 = file_donut_pie_for_month(counts_by_centre(df_for_kwd2),"3rd")

## Add the line graphs and donut pies to the Top 3 Customer Interests chart
logger.info("Adding line graphs, donuts and customers to the Top 3 Customer Interests slide (5th slide)")
s = prs.slides[4]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Top 3 Customer Interests: "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
i = find_text_in_shapes(slide_shapes,"Topic1")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = kwd0
else:
    logger.error("Could not find Topic1 placeholder on 5th slide")
i = find_text_in_shapes(slide_shapes,"Topic2")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = kwd1
else:
    logger.error("Could not find Topic2 placeholder on 5th slide")
i = find_text_in_shapes(slide_shapes,"Topic3")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    text_frame.text = kwd2
else:
    logger.error("Could not find Topic3 placeholder on 5th slide")
#Add the line graphs and the donuts for each of the topics    
top=Inches(1.8); h=Inches(1.0); w=Inches(1.7)
slide_shapes.add_picture(file_linegraph_topic1,Inches(0.7),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic2,Inches(4.7),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic3,Inches(8.8),top,height=h,width=w)
w=Inches(2.0)
slide_shapes.add_picture(file_donut_topic1,Inches(2.5),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic2,Inches(6.55),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic3,Inches(10.6),top,height=h,width=w)
#Find where the placeholders are for the customer lists and update them
i = find_text_in_shapes(slide_shapes,"Customers1")
if (i>=0): 
    logger.debug("Writing list of %i customers for first keyword" % (len(df_for_kwd0)))
    write_customer_list(df_for_kwd0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers1 placeholder")
    
i = find_text_in_shapes(slide_shapes,"Customers2")
if (i>=0): 
    logger.debug("Writing list of",len(df_for_kwd1),"customers for second keyword")
    write_customer_list(df_for_kwd1,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers2 placeholder")
    
i = find_text_in_shapes(slide_shapes,"Customers3")
if (i>=0): 
    logger.debug("Writing list of",len(df_for_kwd2),"customers for third keyword")
    write_customer_list(df_for_kwd2,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers3 placeholder")

################################################
## Now generate the Industry Insights donuts and top keyword lists
################################################
industry_counts = df_for_3months["Updated Industry Sep6"].value_counts()
file_donut_ind_vols = file_donut_pie_for_industries(industry_counts)

df_for_FinSvcs = df_for_3months[df_for_3months["Updated Industry Sep6"]=="Fin Svcs"]
file_donut_FinSvcs = file_donut_pie_for_month(counts_by_centre(df_for_FinSvcs),"FinSvcs")
kwd_counts_for_FinSvcs = keywords_in_dataframe(df_for_FinSvcs,year=yyyy,month=mm)
logger.info("Top keywords and counts for FinSvcs: %r" % (list(kwd_counts_for_FinSvcs.items())[:4]) )

df_for_PublicSector = df_for_3months[df_for_3months["Updated Industry Sep6"]=="Public Sector"]
file_donut_PublicSector = file_donut_pie_for_month(counts_by_centre(df_for_PublicSector),"PublicSector")
kwd_counts_for_PublicSector = keywords_in_dataframe(df_for_PublicSector,year=yyyy,month=mm)
logger.info("Top keywords and counts for Public Sector: %r" % (list(kwd_counts_for_PublicSector.items())[:4]) )

df_for_Mfg = df_for_3months[df_for_3months["Updated Industry Sep6"]=="Mfg"]
file_donut_Mfg = file_donut_pie_for_month(counts_by_centre(df_for_Mfg),"Mfg")
kwd_counts_for_Mfg = keywords_in_dataframe(df_for_Mfg,year=yyyy,month=mm)
logger.info("Top keywords and counts for Manufacturing: %r" % (list(kwd_counts_for_Mfg.items())[:4]) )

df_for_CME = df_for_3months[df_for_3months["Updated Industry Sep6"]=="CME"]
file_donut_CME = file_donut_pie_for_month(counts_by_centre(df_for_CME),"CME")
kwd_counts_for_CME = keywords_in_dataframe(df_for_CME,year=yyyy,month=mm)
logger.info("Top keywords and counts for CME: %r" % (list(kwd_counts_for_CME.items())[:4]) )

df_for_RCG = df_for_3months[df_for_3months["Updated Industry Sep6"]=="RCG"]
file_donut_RCG = file_donut_pie_for_month(counts_by_centre(df_for_RCG),"RCG")
kwd_counts_for_RCG = keywords_in_dataframe(df_for_RCG,year=yyyy,month=mm)
logger.info("Top keywords and counts for RCG: %r" % (list(kwd_counts_for_RCG.items())[:4]) )


    
## Now do the by-industry charts
logger.info("Adding main donut to Industry Insights slide (9th slide)")
s = prs.slides[8]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Industry Insights "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)
#Add the individual industry donuts broken down by centre  
logger.info("Adding donuts for top industries (9th slide)")
left=Inches(1.7); top=Inches(1.9)
slide_shapes.add_picture(file_donut_ind_vols,left,top,height=Inches(4.5),width=Inches(3.7))
left=Inches(8.3); h=Inches(0.9); w=Inches(1.85)
slide_shapes.add_picture(file_donut_FinSvcs,     left,Inches(1.25), height=h,width=w)
slide_shapes.add_picture(file_donut_PublicSector,left,Inches(2.35), height=h,width=w)
slide_shapes.add_picture(file_donut_Mfg,         left,Inches(3.40), height=h,width=w)
slide_shapes.add_picture(file_donut_CME,         left,Inches(4.50), height=h,width=w)
slide_shapes.add_picture(file_donut_RCG,         left,Inches(5.60), height=h,width=w)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
logger.info("Adding top interests for each Industry (9th slide)")
i = find_text_in_shapes(slide_shapes,"Top interests - FS")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    write_top_keywords(text_frame,"Top Interests - FS",kwd_counts_for_FinSvcs,count_rows_with_comments(df_for_FinSvcs))
else:
    logger.error("Could not find <Top interests - FS> placeholder on 9th slide")

i = find_text_in_shapes(slide_shapes,"Top interests - PS")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    write_top_keywords(text_frame,"Top Interests - PS",kwd_counts_for_PublicSector,count_rows_with_comments(df_for_PublicSector))
else:
    logger.error("Could not find <Top interests - PS> placeholder on 9th slide")

i = find_text_in_shapes(slide_shapes,"Top interests - Mfg")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    write_top_keywords(text_frame,"Top Interests - Mfg",kwd_counts_for_Mfg,count_rows_with_comments(df_for_Mfg))
else:
    logger.error("Could not find <Top interests - Mfg> placeholder on 9th slide")
    
i = find_text_in_shapes(slide_shapes,"Top interests - CME")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    write_top_keywords(text_frame,"Top Interests - CME",kwd_counts_for_CME,count_rows_with_comments(df_for_CME))
else:
    logger.error("Could not find <Top interests - CME> placeholder on 9th slide")

i = find_text_in_shapes(slide_shapes,"Top interests - RCG")
if (i>=0): 
    text_frame = slide_shapes[i].text_frame 
    write_top_keywords(text_frame,"Top Interests - RCG",kwd_counts_for_RCG,count_rows_with_comments(df_for_RCG))
else:
    logger.error("Could not find <Top interests - RCG> placeholder on 9th slide")
    
## Close the source presentation
logger.info("Saving Powerpoint file for "+this_month)
prs.save('GCA_Customer_Insights_'+this_month+'-'+str(yyyy)+'.pptx')
## Close any open figures
plt.close("all")
logger.info("...and we're done!")
for h in list(logger.handlers): logger.removeHandler(h)   # may be several here if we've crashed sometimes
