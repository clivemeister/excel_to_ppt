import pandas as pd
from datetime import datetime, date
import matplotlib.pyplot as plt
import collections
from collections import Counter
import calendar
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import os
import sys
import logging
import re

logger = logging.getLogger(__name__)
console=logging.StreamHandler()
console.setLevel(logging.INFO)
formatter=logging.Formatter('%(asctime)s %(levelname)s %(message)s')
console.setFormatter(formatter)
logger.addHandler(console)

tmpdir = "tmp/"
icondir = "icons/"
yyyy,mm=2017,12
logger.info('Starting run for %i-%i...' % (yyyy,mm))

# Read in the ini file
import configparser
import json
ini_file = 'excel_to_ppt.ini'
cfg = configparser.ConfigParser()
cfg.optionxform = str    # read strings as-is from INI file (default is to lowercase keys)
cfg.read(ini_file)

colour_list=[]   # list of the colour code values
if 'colours' in cfg:
    colour_codes = dict(cfg.items('colours'))
    for i in cfg.items('colours'):
        colour_list.append(i[1])       # put this hex value in list of available colours
else:
    logger.error('No [colours] section in {}'.format(ini_file))

if 'keywords' in cfg:
    c_to_k=cfg.items('keywords')
    dict_colour_of_keywords={}
    for i in cfg.items('keywords'):
        if i[0] in colour_codes:
            dict_colour_of_keywords[colour_codes[i[0]]]=json.loads(i[1])
        else:
            logger.error('In {}, section [keywords], found colour {} which was not listed in [colours] section'.format(ini_file,i[0]))
else:
    logger.error('No [keywords] section in {}'.format(ini_file))

#Build a list of all the keywords.  Use all the things we have a colour for as the list items.
vocab=[]
for i in dict_colour_of_keywords.values():
    for j in i:
        vocab.append(j)

# load the list of synonyms (or mis-spellings) of the keywords
synonym_list={}
for i in cfg.items('synonyms'):
    lst=json.loads(i[1])
    for j in lst:
        synonym_list[j]=i[0]

industry_list=[]   # list of the industry code values
if 'industries' in cfg:
    industry_longnames = dict(cfg.items('industries'))
    for i in cfg.items('industries'):
        industry_list.append(i[0])
else:
    logger.error('No [industries] section in {}'.format(ini_file))

centres = ["H","NY1","SNG","LON1","PA"]   # short names used in the Excel file
centres_long = ["Houston", "New York", "Singapore", "London", "Palo Alto"]  # same order as short names

stop_words=[]
if 'stopwords' in cfg:
    stop_words+=json.loads(cfg.get('stopwords','stop_words'))
else:
    logger.error('No [stopwords] section in {}'.format(ini_file))
stop_words+=vocab

def make_date(cell_val):
    """Used to create a better-structured date value from the sometimes-odd values in the date column
       Returns a datetime
    """
    if type(cell_val) is datetime:
        v=cell_val.date()
    elif type(cell_val) is str:
        try:
            v=datetime.strptime(cell_val,"%b %d, %Y").date()
        except ValueError:
            v=date.fromordinal(1)
    else:
        v=date.fromordinal(2)
    return v

def tidy_text(cell_val):
    """Standardises the text in a cell:
       Returns the tidied text
    """
    if type(cell_val) is str:
        cell = cell_val.lower()
        cell=re.sub('[\n&@,.:-]',' ',cell)
        cell=" ".join(cell.split())   # idiom to turn multiple spaces between words into single spaces
        for k,v in synonym_list.items():
            cell=cell.replace(k,v)
    else:
        cell=""
    return cell

def replace_strings(series,repl_dict):
    for k,v in repl_dict.items():
        series=series.str.replace(k,v)
    return(series)

def keyword_counts_in_series(series,keywords):
    list_of_strings = series
    if keywords is str:
        #get here only if keywords is a single string, not a list
        count_list=sum(list_of_strings.str.find(keywords)>=0)
    else:
        count_list=Counter()
        for i in keywords:
            count_list[i]=sum(list_of_strings.str.find(i)>=0)
    return count_list

def print_new_candidate_words(df,stop_words,top_n=25):
    """Print the top_n words in the relevant columns in the dataframe df that aren't in the
       stop_words
       Input are a dataframe with columns 'wtlma' and 'ai',
       a list in stop_words of single words to ignore, and the number of top words to show (default 25)
    """
    big_lst=[]
    for lst in list(df.wtlma.str.split()):
        if type(lst)==list: big_lst += lst
    for lst in list(df.ai.str.split()):
        if type(lst)==list: big_lst += lst

    all_ctr=Counter(big_lst)
    for s in stop_words: del all_ctr[s]
    print("Common words which are possible candidates for new keywords:")
    print(all_ctr.most_common(20))
    return

def dataframe_for_month(df, year=2017, month=1):
    """Yields a subset the dataframe with only those rows in the given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    if (mm==12): mm2,yyyy2=1,yyyy+1
    else: mm2,yyyy2=mm+1,yyyy

    logger.debug("Selecting rows for %i-%i" % (yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy,mm,1)) & (df.date<date(yyyy2,mm2,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def dataframe_for_6months(df, year=2017, month=1):
    """Yields a subset the dataframe with those rows for last 6 months up to given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    # calc month after this one, to set upper limit for search
    if (mm==12): mm_end,yyyy_end=1,yyyy+1
    else: mm_end,yyyy_end=mm+1,yyyy
    # calc 6 months ago, to set lower limit for search
    if (mm<=6): mm_start,yyyy_start=mm+6,yyyy-1
    else: mm_start,yyyy_start=mm-6,yyyy

    logger.debug("Selecting rows for %i-%i to %i-%i" % (yyyy_start,mm_start,yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy_start,mm_start,1)) & (df.date<date(yyyy_end,mm_end,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def found_word_list(df, kwd):
    """Passed a dataframe with columns of at least 'wtlma' and 'ai', together with a keyword
       Return a list (actually a Series) with True wherever the dataframe contains kwd or a synonym in the relevant columns,
       and False otherwise.  Can be used to subset the original dataframe to pick out the rows with the keyword in them:
           df[found_word_list(df,"foo")]   => subset of df with "foo" in one of the columns
    """
    strings_wtlma = replace_strings(df.wtlma, synonym_list)
    strings_ai = replace_strings(df.ai, synonym_list)
    found_list = (strings_wtlma.str.find(kwd)>=0) | (strings_ai.str.find(kwd)>=0)
    return found_list

def keywords_in_dataframe(df,keyword_list):
    """ Count each keyword in the dataframe's important columns
        Returns a Counter collection of (keyword:count)
    """
    import operator   #for sorted(), used in sorting items into OrderedDict

    logger.debug("Counting occurrences of keywords for this month")
    counter_words_to_freq = Counter()
    for w in keyword_list:
         counter_words_to_freq[w] += sum(found_word_list(df,w))

    return counter_words_to_freq

def count_rows_with_comments(df):
    """ Count the number of rows in dataframe with a comment in either <Want to Learn More About> or <Action Items>
        Returns the count
    """
    return (df["Want to Learn More About"].notnull() | df["Action Items"].notnull()).sum()

def counts_by_centre(dataframe):
    centre_counts = dataframe.Ctr.value_counts()
    counts_list = (centre_counts.PA if hasattr(centre_counts,"PA") else 0,
                   centre_counts.H if hasattr(centre_counts,"H") else 0,
                   centre_counts.NY1 if hasattr(centre_counts,"NY1") else 0,
                   centre_counts.LON1 if hasattr(centre_counts,"LON1") else 0,
                   centre_counts.SNG if hasattr(centre_counts,"SNG") else 0
                   )
    return counts_list

class SimpleGroupedColorFunc(object):
    """Create a color function object which assigns EXACT colors
       to certain words based on the color to words mapping

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        self.word_to_color = {word: color
                              for (color, words) in color_to_words.items()
                              for word in words}

        self.default_color = default_color

    def __call__(self, word, **kwargs):
        return self.word_to_color.get(word, self.default_color)

class GroupedColorFunc(object):
    """Create a color function object which assigns DIFFERENT SHADES of
       specified colors to certain words based on the color to words mapping.

       Uses wordcloud.get_single_color_func

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        from wordcloud import get_single_color_func
        self.color_func_to_words = [
            (get_single_color_func(color), set(words))
            for (color, words) in color_to_words.items()]

        self.default_color_func = get_single_color_func(default_color)

    def get_color_func(self, word):
        """Returns a single_color_func associated with the word"""
        try:
            color_func = next(
                color_func for (color_func, words) in self.color_func_to_words
                if word in words)
        except StopIteration:
            color_func = self.default_color_func

        return color_func

    def __call__(self, word, **kwargs):
        return self.get_color_func(word)(word, **kwargs)

def file_wordcloud_for_month(keywords_for_month, useful_rows_for_month, year=2017, month=8):
    """Input is a keywords_for_month, a Counter, plus useful_rows_for_month, an integer, and year and month as integers
       Returns the filename where the wordcloud is stored
    """
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    # (to install, run "pip install wordcloud"
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt
    import os

    # set font so 30% occurrence of top word uses 196 point font
    percent = round(100*(keywords_for_month.most_common(1)[0][1]/useful_rows_for_month))
    font_for_biggest_word = round( 196 * percent/30 )
    logger.debug("font for word {} is {} based on {}% = {} / {}".format(
                               keywords_for_month.most_common(1)[0][0],
                                     font_for_biggest_word,
                                             percent,
                                                keywords_for_month.most_common(1)[0][1],
                                                        useful_rows_for_month))

    font_path = os.path.join("C:",os.sep,"Windows",os.sep,"Fonts",os.sep,'arial.ttf')
    logger.debug("Generating the wordcloud")
    keywords_for_month += Counter()    # remove any zero or negative counts from the list
    wc_for_month = WordCloud(font_path=font_path,
                             width=2500,height=500,
                             prefer_horizontal=1.0,
                             relative_scaling=0.7,
                             max_font_size=font_for_biggest_word,
                             background_color="white",
                             random_state=1
                             ).generate_from_frequencies(keywords_for_month)


    # Words that are not in any of the dict_colour_of_keywords values
    # will be colored with a grey single color function
    default_color = 'grey'

    # Apply our color function
    #grouped_color_func = GroupedColorFunc(dict_colour_of_keywords, default_color)
    grouped_color_func = SimpleGroupedColorFunc(dict_colour_of_keywords, default_color)
    wc_for_month.recolor(color_func=grouped_color_func)

    # Build the generated image
    plt.imshow(wc_for_month,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    #plt.show()

    filename = tmpdir+"wordcloud-"+calendar.month_name[month]+".png"
    plt.imsave(filename,wc_for_month,format="png")
    plt.close()

    return filename


def file_graph_for_month_kwd(kwd,kwd_pos,vals,months,line_color):
    fig,ax=plt.subplots(figsize=(5.75,3.25))
    ##############ax.clear()

    ax.plot(vals,line_color,linewidth=3)
    ax.set_axis_off()
    plt.ylim(min(vals)-0.05,max(vals)+0.05)
    m_minus_2_percent = "{0:.0f}%".format(vals[0] * 100)
    m_this_percent    = "{0:.0f}%".format(vals[2] * 100)
    ax.text(0,vals[0],calendar.month_name[months[0]]+"\n"+m_minus_2_percent,fontsize=36)
    ax.text(2,vals[2],calendar.month_name[months[2]]+"\n"+m_this_percent,fontsize=36)
    filename = tmpdir+"graph-"+str(kwd_pos)+".png"
    logger.debug("Saving %s graph for keyword %s in file %s" % (kwd_pos,kwd,filename))
    fig.savefig(filename,bbox_inches="tight")
    plt.close(fig)

    return filename

def file_donut_pie_for_month(values,name):
    fig, ax = plt.subplots()
    ax.axis('equal')
    outside, _ = ax.pie(values,startangle=90,counterclock=False,
                        colors=list(colour_list))
    ax.legend(("Palo Alto","Houston","NY","London","Singapore"),fontsize=24,bbox_to_anchor=(0.8,1.0),frameon=False)
    width = 0.50  #determines the thickness of donut rim
    plt.setp( outside, width=width, edgecolor='white')

    filename = tmpdir+"donut-"+re.sub(r"[& ]","_",str(name))+".png"
    logger.debug("Saving <%s> donut in file %s with values %r" % (name, filename, values))
    fig.savefig(filename,bbox_inches="tight")

    plt.close(fig)

    return filename

def file_donut_pie_for_industries(industries):
    fig, ax = plt.subplots()
    ax.axis('equal')
    width = 2.0
    outside, _ = ax.pie(industries.values,
                        startangle=90,counterclock=False,
                        colors=list(colour_list),
                        radius=5.0)
    ax.legend(industries.index.tolist(),ncol=3,fontsize=24,loc=10,bbox_to_anchor=(0.5,-2),frameon=False)
    plt.setp( outside, width=width, edgecolor='white')

    filename = tmpdir+"donut-industries.png"
    logger.debug("Saving Industries donut in file %s with values %r" % (filename, industries))
    fig.savefig(filename,bbox_inches="tight")

    plt.close(fig)

    return filename

def write_customer_list(df_for_kwd,text_frame):
    text_frame.text = "Briefings"
    font=text_frame.paragraphs[0].font
    font.size = Pt(12)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for i in (df_for_kwd)["Account Name"].str.title().iteritems():
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        new_para.text=i[1]

    return

def write_top_keywords(text_frame,titleText,kwd_counts_for_ind,rowcount,percent=True):
    """Write a section headed by titleText, with bullets of the top 4 elements of Counter
       kwd_counts_for_ind plus their percentage of use (if percent=True), where
       the percentage is calculated with rowcount as the denominator.
         e.g:
           c=Counter([("alpha":5),("bravo":4),("charlie":3),("delta":2))])
           write_top_keywords(tf,"My list",c,10)
         yields:
           My list
            - alpha - 50%
            - bravo - 40%
            - charlie - 30%
    """
    text_frame.text = titleText
    font=text_frame.paragraphs[0].font
    font.size = Pt(10)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for k,v in kwd_counts_for_ind.most_common(4):
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        if (percent==True):
            new_para.text=" - "+k+" - "+"{0:.0f}%".format(v/rowcount * 100)
        else:
            new_para.text=" - "+k

    return

def find_text_in_shapes(slide_shapes,searchword):
    """Find the given text in the list of powerpoint shapes passed in slide_shapes
    """
    found_idx = -1
    for i in range(len(slide_shapes)):
        if (slide_shapes[i].has_text_frame):
            text_frame = slide_shapes[i].text_frame
            if (text_frame.text.find(searchword)>=0):
                found_idx = i
                break

    return found_idx

def replace_text_in_shape(slide_shapes,find,use,slidename):
    """Find the placeholder text "find" in the slide_shapes, replace it with "use",
       reporting an error on "slidename" if we can't find "find"
    """
    i = find_text_in_shapes(slide_shapes,find)
    if (i>=0):
        text_frame = slide_shapes[i].text_frame
        text_frame.text = use
    else:
        logger.error("Could not find %s placeholder on %s" % (find,slidename))

    return

def new_run_in_slide(para,text='text',fontname='Arial',fontsize=24):
    "Add new run to a paragraph in a powerpoint slide"
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.name = fontname
    return run

import struct
import imghdr

def get_image_size(fname):
    '''Determine the image type of fname and return its width, height.
    from draco (via stackoverflow)'''
    with open(fname, 'rb') as fhandle:
        head = fhandle.read(24)
        if len(head) != 24:
            return
        if imghdr.what(fname) == 'png':
            check = struct.unpack('>i', head[4:8])[0]
            if check != 0x0d0a1a0a:
                return
            width, height = struct.unpack('>ii', head[16:24])
        elif imghdr.what(fname) == 'gif':
            width, height = struct.unpack('<HH', head[6:10])
        elif imghdr.what(fname) == 'jpeg':
            try:
                fhandle.seek(0) # Read 0xff next
                size = 2
                ftype = 0
                while not 0xc0 <= ftype <= 0xcf:
                    fhandle.seek(size, 1)
                    byte = fhandle.read(1)
                    while ord(byte) == 0xff:
                        byte = fhandle.read(1)
                    ftype = ord(byte)
                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
                # We are at a SOFn block
                fhandle.seek(1, 1)  # Skip `precision' byte.
                height, width = struct.unpack('>HH', fhandle.read(4))
            except Exception: #IGNORE:W0703
                return
        else:
            return
        return width, height

def add_icon(ss,icon_name,left,top,small=False):
    ## ok if icon exists, and use error icon if not
    logger.debug("adding icon {}".format(icon_name))
    try:
        filename = icondir+icon_name+".png"
        icon_width, icon_height = get_image_size(filename)
        pixels_per_mm = 6
        if small:
            icon_width = 0.4*icon_width
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(12),Mm(12))
        else:
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(34),Mm(34))
    except:
        logger.warning("Could not find icon {}".format(icondir+icon_name+'.png'))
    return

excel_file='Insights.xlsx'
logger.info("Importing excel file "+excel_file)
all_df = pd.read_excel(open(excel_file,'rb'),header=8,usecols="A:S")

logger.debug("Adding structured columns")
all_df.insert(loc=0,column='date',value=all_df['Visit Date'].apply(make_date))
all_df.insert(loc=1,column='wtlma',value=all_df['Want to Learn More About'].apply(tidy_text))
all_df.insert(loc=2,column='ai',value=all_df['Action Items'].apply(tidy_text))

print_new_candidate_words(all_df,stop_words,top_n=40)

## Given the month and year, calc the number of the previous few months
if (mm==1): mm_minus_1,year_for_mm_minus_1 = 12, yyyy-1
else: mm_minus_1, year_for_mm_minus_1 = mm-1, yyyy
if (mm<=2): mm_minus_2,year_for_mm_minus_2 = mm+10, yyyy-1
else: mm_minus_2, year_for_mm_minus_2 = mm-2, yyyy
if (mm<=6): mm_minus_6,year_for_mm_minus_6 = mm+6, yyyy-1
else: mm_minus_6, year_for_mm_minus_6 = mm-6, yyyy

### Now start to generate the powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt, Mm

# These appear to be the layouts used in the master for this slide deck
LAYOUT_TITLE_WITH_DARK_PICTURE = 0
LAYOUT_TITLE_SLIDE_WITH_NAME   = 1
LAYOUT_DIVIDER                 = 2
LAYOUT_TITLE_ONLY              = 3
LAYOUT_TITLE_AND_SUBTITLE      = 4
LAYOUT_BLANK                   = 5

## Open up the source presentation
prs = Presentation('GCA_Customer_Insights_Month-Year.pptx')
this_month = calendar.month_name[mm]
earliest_month = calendar.month_name[mm_minus_2]
logger.info("Creating presentation for %s" % (this_month))

################################################
## Modify existing title slide with month
################################################
logger.info("Modifying text in slide 0")
s = prs.slides[0]
slide_shapes = s.shapes
text_frame = slide_shapes[0].text_frame  # should be the the title textframe
#clear existing text, and write new text into title textframe
text_frame.clear()
# First para with a run of 60-point Arial font
new_run_in_slide(text_frame.paragraphs[0],text='Customer Insights',fontname='Arial',fontsize=60)
# Second para with a run of 32-point font
new_run_in_slide(text_frame.add_paragraph(),text='Learnings from '+this_month+' EBC/CEC visits',fontsize=32)

################################################
## 3rd slide: count the keywords for this month and build the wordcloud
################################################
logger.info("3rd slide: large wordcloud for this month, and top 3 keywords")
df_for_month = dataframe_for_month(all_df, year=yyyy, month=mm)
kwd_count_for_month = keywords_in_dataframe(df_for_month,vocab)
useful_rows_in_m = count_rows_with_comments(df_for_month)
logger.debug("Top keyword/counts for month %i : %r" % (mm,kwd_count_for_month.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_month, useful_rows_in_m,
                         year=yyyy,month=mm)
## Modifying main wordcloud slide by changing title and adding pic for this month's wordcloud
logger.debug("Modifying text and adding wordcloud in slide 2")
s = prs.slides[2]
slide_shapes=s.shapes
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="In "+this_month+" customers wanted to learn more about...",
       fontname="Arial",fontsize=28)
left=Mm(12.5); top=Mm(100)
slide_shapes.add_picture(tmpdir+"wordcloud-"+this_month+".png",left,top,height=Mm(63))

##Find where the placeholders are for the top 3 keywords update them
top_3 = kwd_count_for_month.most_common(3)   # top 3 keywords for most recent month in list with their counts
kwd0=top_3[0][0]
kwd1=top_3[1][0]
kwd2=top_3[2][0]
logger.debug("Top 3 keywords for current month are %s %s %s" % (kwd0,kwd1,kwd2))
replace_text_in_shape(slide_shapes,find="topword_1",use=kwd0,slidename="3rd slide")
replace_text_in_shape(slide_shapes,find="topword_2",use=kwd1,slidename="3rd slide")
replace_text_in_shape(slide_shapes,find="topword_3",use=kwd2,slidename="3rd slide")
t=Mm(51)
add_icon(slide_shapes,kwd0,left=Mm(47),top=t)
add_icon(slide_shapes,kwd1,left=Mm(152),top=t)
add_icon(slide_shapes,kwd2,left=Mm(255),top=t)

################################################
## 4th slide: count the keywords for previous two months, and build their wordclouds.
################################################
logger.info(">>>> 4th slide: three wordclouds for most recent 3 months")
df_for_month_minus_1 = dataframe_for_month(all_df, year=year_for_mm_minus_1, month=mm_minus_1)
kwd_count_for_m_minus_1 = keywords_in_dataframe(df_for_month_minus_1,vocab)
logger.debug("Top keyword/counts for month %i : %r" % (mm_minus_1,kwd_count_for_m_minus_1.most_common(5)) )

df_for_month_minus_2 = dataframe_for_month(all_df, year=year_for_mm_minus_2, month=mm_minus_2)
kwd_count_for_m_minus_2 = keywords_in_dataframe(df_for_month_minus_2,vocab)
logger.debug("Top keyword/counts for month %i : %r" % (mm_minus_2,kwd_count_for_m_minus_2.most_common(5)) )

useful_rows_in_m_2 = count_rows_with_comments(df_for_month_minus_2)
useful_rows_in_m_1 = count_rows_with_comments(df_for_month_minus_1)
logger.debug("Useful rows in months 2,1,0 are %i, %i, %i" % (useful_rows_in_m_2,useful_rows_in_m_1,useful_rows_in_m))

file_wordcloud_for_month(kwd_count_for_m_minus_1, useful_rows_in_m_1, year=year_for_mm_minus_1,month=mm_minus_1)
file_wordcloud_for_month(kwd_count_for_m_minus_2, useful_rows_in_m_2, year=year_for_mm_minus_2,month=mm_minus_2)

## Modifying 3-month wordcloud slide by adding pic for last 3 months' wordcloud
logger.info("Adding three wordclouds in 4th slide")
s = prs.slides[3]
slide_shapes=s.shapes
left=Mm(93)
top=Mm(36)
h=Mm(38)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_2]+".png",
                         left,top,height=h)
top=Mm(83)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_1]+".png",
                         left,top,height=h)
top=Mm(128)
slide_shapes.add_picture(tmpdir+"wordcloud-"+this_month+".png",
                         left,top,height=h)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
replace_text_in_shape(slide_shapes,find="Month-2",use=calendar.month_name[mm_minus_2],slidename="4th slide")
replace_text_in_shape(slide_shapes,find="Month-1",use=calendar.month_name[mm_minus_1],slidename="4th slide")
replace_text_in_shape(slide_shapes,find="Month-0",use=calendar.month_name[mm],slidename="4th slide")

################################################
## 5th slide: for top 3 keywords for this month graph their usage
################################################
logger.info(">>>> 5th slide: top keywords for this month")
months=[mm_minus_2,mm_minus_1,mm]

kwd0_c2 = kwd_count_for_m_minus_2[kwd0] if (kwd0 in kwd_count_for_m_minus_2) else 0
kwd0_c1 = kwd_count_for_m_minus_1[kwd0] if (kwd0 in kwd_count_for_m_minus_1) else 0
kwd0_c0 = kwd_count_for_month[kwd0]     # must have keyword as it came from this dictionary
vals_kwd0=[kwd0_c2/useful_rows_in_m_2, kwd0_c1/useful_rows_in_m_1, kwd0_c0/useful_rows_in_m]
file_linegraph_topic1 = file_graph_for_month_kwd(kwd0,"1st",vals_kwd0,months,colour_list[0])
logger.debug("Kwd0 is %s, data %r" % (kwd0,vals_kwd0))
kwd1_c2 = kwd_count_for_m_minus_2[kwd1] if (kwd1 in kwd_count_for_m_minus_2) else 0

kwd1_c1 = kwd_count_for_m_minus_1[kwd1] if (kwd1 in kwd_count_for_m_minus_1) else 0
kwd1_c0 = kwd_count_for_month[kwd1]     # must have keyword as it came from this dictionary
vals_kwd1=[kwd1_c2/useful_rows_in_m_2, kwd1_c1/useful_rows_in_m_1, kwd1_c0/useful_rows_in_m]
file_linegraph_topic2 = file_graph_for_month_kwd(kwd1,"2nd",vals_kwd1,months,colour_list[1])
logger.debug("Kwd1 is %s, data %r" % (kwd1,vals_kwd1))

kwd2_c2 = kwd_count_for_m_minus_2[kwd2] if (kwd2 in kwd_count_for_m_minus_2) else 0
kwd2_c1 = kwd_count_for_m_minus_1[kwd2] if (kwd2 in kwd_count_for_m_minus_1) else 0
kwd2_c0 = kwd_count_for_month[kwd2]     # must have keyword as it came from this dictionary
vals_kwd2=[kwd2_c2/useful_rows_in_m_2, kwd2_c1/useful_rows_in_m_1, kwd2_c0/useful_rows_in_m]
file_linegraph_topic3 = file_graph_for_month_kwd(kwd2,"3rd",vals_kwd2,months,colour_list[2])
logger.debug("Kwd0 is %s, data %r" % (kwd2,vals_kwd2))

## Build a subset of the dataframe for last 3 months that uses each of the top 3 kwds in this month
df_for_3months = pd.concat([df_for_month,df_for_month_minus_1,df_for_month_minus_2])
df_for_kwd0 = df_for_3months.loc[found_word_list(df_for_3months,kwd0)]
df_for_kwd1 = df_for_3months.loc[found_word_list(df_for_3months,kwd1)]
df_for_kwd2 = df_for_3months.loc[found_word_list(df_for_3months,kwd2)]

## Create donut pies showing split of visits expressing interest in top 3 topics by centre over last 3 months
kwd0_counts = counts_by_centre(df_for_kwd0)
file_donut_topic1 = file_donut_pie_for_month(counts_by_centre(df_for_kwd0),"1st")
file_donut_topic2 = file_donut_pie_for_month(counts_by_centre(df_for_kwd1),"2nd")
file_donut_topic3 = file_donut_pie_for_month(counts_by_centre(df_for_kwd2),"3rd")

## Add the line graphs and donut pies to the Top 3 Customer Interests chart
logger.info("Adding line graphs, donuts and customers to the Top 3 Customer Interests slide (5th slide)")
s = prs.slides[4]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Top 3 Customer Interests: "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
replace_text_in_shape(slide_shapes,find="Topic1",use=kwd0,slidename="5th slide")
replace_text_in_shape(slide_shapes,find="Topic2",use=kwd1,slidename="5th slide")
replace_text_in_shape(slide_shapes,find="Topic3",use=kwd2,slidename="5th slide")

#Add the line graphs and the donuts for each of the topics
top=Mm(46); h=Mm(25); w=Mm(43)
slide_shapes.add_picture(file_linegraph_topic1,Mm(18),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic2,Mm(120),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic3,Mm(220),top,height=h,width=w)
top=Mm(43); h=Mm(28); w=Mm(57);
slide_shapes.add_picture(file_donut_topic1,Mm(58),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic2,Mm(160),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic3,Mm(263),top,height=h,width=w)

#Find where the placeholders are for the customer lists and update them
df_for_kwd0_month0 = df_for_month.loc[found_word_list(df_for_month,kwd0)]
df_for_kwd1_month0 = df_for_month.loc[found_word_list(df_for_month,kwd1)]
df_for_kwd2_month0 = df_for_month.loc[found_word_list(df_for_month,kwd2)]
i = find_text_in_shapes(slide_shapes,"Customers1")
if (i>=0):
    logger.debug("Writing list of %i customers for first keyword" % (len(df_for_kwd0)))
    write_customer_list(df_for_kwd0_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers1 placeholder")

i = find_text_in_shapes(slide_shapes,"Customers2")
if (i>=0):
    logger.debug("Writing list of",len(df_for_kwd1),"customers for second keyword")
    write_customer_list(df_for_kwd1_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers2 placeholder")

i = find_text_in_shapes(slide_shapes,"Customers3")
if (i>=0):
    logger.debug("Writing list of",len(df_for_kwd2),"customers for third keyword")
    write_customer_list(df_for_kwd2_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers3 placeholder")

################################################
## 9th slide: Now generate the Industry Insights donuts and top keyword lists
################################################
logger.info(">>>> 9th slide: industries and how they show up across the centres")
## Generate the image for the big donut showing volume for all industries across the months
industry_counts = df_for_3months["Updated Industry Sep6"].value_counts()
file_donut_ind_vols = file_donut_pie_for_industries(industry_counts)

## Generate the individual donuts showing, for each industry, the spread across the centres


# set up some dictionaries to hold the industry datasets for later
df_for_ind={};
file_donut_for_ind={};
kwd_counts_for_ind={}

for ind in industry_list:
    df_for_ind[ind] = df_for_3months[df_for_3months["Updated Industry Sep6"]==ind]
    file_donut_for_ind[ind] = file_donut_pie_for_month(counts_by_centre(df_for_ind[ind]),ind)
    kwd_counts_for_ind[ind] = keywords_in_dataframe(df_for_ind[ind],vocab)
    logger.debug("Top keyword/counts for %s: %r" % (ind,kwd_counts_for_ind[ind].most_common(4)) )

## Now write out the charts and text on industry (9th) slide
s = prs.slides[8]
slide_shapes=s.shapes

#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Industry Insights "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)

#Add the one big donut showing volumes for each industry to the slide
logger.info("Adding main donut to Industry Insights slide (9th slide)")
left=Mm(43); top=Mm(48)
slide_shapes.add_picture(file_donut_ind_vols,left,top,height=Mm(115),width=Mm(94))

#Add the individual industry donuts broken down by centre to the slide
#Need to pick the top 5, excluding "Other", from the list in industry_list
logger.info("Adding donuts for top industries in each centre (9th slide)")
left=Mm(208); top=(Mm(30),Mm(58),Mm(86),Mm(115),Mm(142))
h=Mm(25); w=Mm(52)
n=0 # used to count the number of industry pies we have placed (we can't use enumerate(industry_counts) as we don't always place a pie)
for ind in industry_counts.index:
    if   (ind!="Other"):
        logger.debug("Writing %s as industry %i" %(industry_longnames[ind],n))
        #Write the industry name as the main title for this box
        replace_text_in_shape(slide_shapes,"Industry-{}".format(n),industry_longnames[ind],"9th slide")
        #Add the donut showing breakdown of centres that hosted this industry
        slide_shapes.add_picture(file_donut_for_ind[ind], left,top[n], height=h,width=w)
        #Now write the list of top interests for this industry
        idx = find_text_in_shapes(slide_shapes,"Top interests - {}".format(n))
        if (idx>=0):
            text_frame = slide_shapes[idx].text_frame
            write_top_keywords(text_frame,"Top Interests - "+industry_longnames[ind],
                               kwd_counts_for_ind[ind],
                               count_rows_with_comments(df_for_ind[ind]),
                               percent=False
                              )
        else:
            logger.error("Could not find <Top interests - {}> placeholder on 9th slide".format(n))
        n+=1
    if (n>=5): break    # Stop after we've put 5 pictures in place

################################################
## 10th slide: Partner insights
################################################
logger.info(">>>> 10th slide: top partner keywords and partner attendance broken out by centre")
# rather copmlex expression to find which rows are partners who are attending as partners, or partner-led briefings for customers
df_for_partners = df_for_3months.loc[ ( df_for_3months['Account Type'].isin(["Channel/ Reseller","Systems Integrator"])
                                        & (df_for_3months['Partner / Customer']!="Customer")
                                      ) | (df_for_3months['Partner / Customer']=="Partner")
                                    ]
kwd_count_for_partners = keywords_in_dataframe(df_for_partners,vocab)
useful_rows_in_partners = count_rows_with_comments(df_for_partners)

df_channel = df_for_partners[ df_for_partners['Account Type']=="Channel/ Reseller" ]
df_SI = df_for_partners[ df_for_partners['Account Type']=="Systems Integrator" ]
df_accompanied = df_for_partners[ df_for_partners['Partner / Customer']=="Partner" ]


SI_count=[len(df_SI[df_SI.Ctr==c]) for c in centres]
Channel_count=[len(df_channel[df_channel.Ctr==c]) for c in centres]
Attended_count=[len(df_accompanied[df_accompanied.Ctr==c]) for c in centres]
logger.debug("Centres being looked at: %r" %(centres))
logger.debug("Volume by centre - partner: %r" %(Channel_count))
logger.debug("Volume by centre - SI: %r" %(SI_count))
logger.debug("Volume by centre - attended: %r" %(Attended_count))

fig, ax = plt.subplots()
fig.set_size_inches(5.5,2.2)
y=(0.5,1,1.5,2,2.5)
ax.barh(y,Channel_count, height=0.35, color=colour_list[0],tick_label=centres_long)
ax.barh(y,SI_count, height=0.35, left=Channel_count,color=colour_list[1])
new_base=[x+y for x,y in zip(Channel_count, SI_count)]
ax.barh(y,Attended_count, height=0.35, left=new_base,color=colour_list[2])
ax.legend(["Channel/ Reseller","SI","Partner attended with customer"],loc='lower center',ncol=3,bbox_to_anchor=(0.5,-0.4))
ax.get_xaxis().set_visible(False)
file_hbar = "barh-PartnerVolumes.png"
fig.savefig(file_hbar,bbox_inches="tight")

plt.close(fig)

## Start to update the slide
s = prs.slides[9]
slide_shapes=s.shapes

#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Partner Insights ("+earliest_month+"-"+this_month+")",
       fontname="Arial",fontsize=28)

# Update the top 4 most common keywords, and their percentages
logger.debug("top 4 partner interests: %r" % (kwd_count_for_partners.most_common(4)))
t=Mm(68)
l=[Mm(30),Mm(68),Mm(103),Mm(140)]
for n,p in enumerate(kwd_count_for_partners.most_common(4)):
    # p is (keyword: count) for each of the top 4 most common keywords
    replace_text_in_shape(slide_shapes,"Interest-{}".format(n),p[0],"10th slide")
    add_icon(slide_shapes,p[0],top=t,left=l[n],small=True)
    replace_text_in_shape(slide_shapes,"Score-{}".format(n),"{0:.0f}%".format(100*p[1]/useful_rows_in_partners),"10th slide")

# Place the horizontal bar graph
slide_shapes.add_picture(file_hbar, Mm(23),Mm(112), height=Mm(56),width=Mm(140))

################################################
## 11th slide: Top interests and industries for last 6 months
################################################
logger.info(">>>> 11th slide: top 5 interests, top 3 industries, and their top interests, by centre, for last 6 months")
df_6months = dataframe_for_6months(all_df, year=yyyy, month=mm)

#build a dictionary of dataframes subsetted by centre, a dictionary of keywords by centre,
#and a dict of top 3 industries per centre and their keywords (where the key is a tuple of (centre, industry) )
dfs_6m_ctr={}
kwd_counts_6m_ctr={}
industry_counts_6m = df_6months["Updated Industry Sep6"].value_counts()
kwd_counts_6m_for_top_inds={}
commented_rows_for_6m={}

for c in centres:
    logger.info("Working on counts for {}".format(c))
    dfs_6m_ctr[c]=df_6months[df_6months.Ctr==c]
    #First, the top keywords for that Centre
    kwd_counts_6m_ctr[c]=keywords_in_dataframe(dfs_6m_ctr[c],vocab)
    logger.debug("Top keyword/counts %s: %r" % (c,kwd_counts_6m_ctr[c].most_common(5)) )
    #Now the top keywords in each industry for that centre
    industry_counts_6m[c] = (dfs_6m_ctr[c])["Updated Industry Sep6"].value_counts()
    for ind in (industry_counts_6m[c]).index:
        this_df = dfs_6m_ctr[c][ dfs_6m_ctr[c]["Updated Industry Sep6"]==ind ]
        kwd_counts_6m_for_top_inds[(c,ind)]=keywords_in_dataframe(this_df, vocab)
        commented_rows_for_6m[(c,ind)]=count_rows_with_comments(this_df)
        logger.debug("Top keyword/counts in %s for %s: %r" % (c,ind,kwd_counts_6m_for_top_inds[(c,ind)].most_common(3)) )


## Build the slide: add the interests, industries, and per-industry interests, for each of the centres
logger.info("Setting the title ")
s = prs.slides[10]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Breakdown by centre ("+calendar.month_name[mm_minus_6]+"-"+this_month+")",
       fontname="Arial",fontsize=28)

#Update, by centre, the top interests, and the top industries with their top interests
logger.info("Setting the per-centre top 5 interests, together with per-centre top 3 industries and their interests")
top_pos=[Mm(48),Mm(73),Mm(102),Mm(127),Mm(152)]   #distances to top of icon for each row
left_pos=[Mm(35),Mm(65),Mm(95),Mm(125),Mm(154)]   #distances to centre of icon for each row
for row,ctr in enumerate(["PA","H","NY1","LON1","SNG"]):  #iterate the centres in the order they appear on the slide
    #First, the top interests for this centre
    for col,p in enumerate(kwd_counts_6m_ctr[ctr].most_common(5)):
        # p is (keyword: count) for each of the keywords, so p[0] is the keyword itself
        replace_text_in_shape(slide_shapes,"{}-interest-{}".format(ctr,col),p[0],"11th slide")
        add_icon(slide_shapes,p[0],left=left_pos[col],top=top_pos[row],small=True)
    #Next, the top industries with their interests for this centre - industry_counts_6m[c] is already ordered highest->lowest count
    n=0  #count how many displayed - need to do this separately from the loop count, as we ignore "Other" as an industry group
    for ind in industry_counts_6m[ctr].index:
        if   (ind!="Other"):
            logger.debug("For centre <%s> industry <%i> is <%s>" %(ctr,n,ind))
            #Write the list of top interests for this industry
            idx = find_text_in_shapes(slide_shapes,"{}-industry-{}".format(ctr,n))
            if (idx>=0):
                write_top_keywords(slide_shapes[idx].text_frame,
                                   ind,
                                   kwd_counts_6m_for_top_inds[(ctr,ind)],
                                   commented_rows_for_6m[(ctr,ind)]
                                  )
            else:
                logger.error("Could not find <{}-industry-{}> placeholder on 11th slide - idx={}".format(ctr,n,idx))
            n+=1   #increment the number of industries written out for this centre
        if (n>=3): break   #exit the loop after writing in 3 industries+interests


################################################
## Close the source presentations
################################################
logger.info("Saving Powerpoint file for "+this_month)
prs.save('GCA_Customer_Insights_'+this_month+'-'+str(yyyy)+'.pptx')
## Close any open figures
plt.close("all")
logger.info("...and we're done!")
for h in list(logger.handlers): logger.removeHandler(h)   # may be several here if we've crashed sometimes