import pandas as pd
from datetime import datetime, date
import matplotlib.pyplot as plt
import collections
import calendar

def make_date(cell_val):
    if type(cell_val) is datetime:
        v=cell_val.date()
    elif type(cell_val) is str:
        try:
            v= datetime.strptime(cell_val,"%b %d, %Y").date()
        except ValueError:
            v=date.fromordinal(1)
    else:
        v=date.fromordinal(2)
    return v

def replace_strings(series,repl_dict):
    for k,v in repl_dict.items():
        series.str.replace(k,v)
    return(series)

def count_strings(series,s):
    if s is str:
        counts=sum(series.str.find(str)>0)
    else:
        counts={}
        for i in s:
            counts[i]=sum(series.str.find(i)>0)
    return counts

def sum_of_dicts(d1,d2):
    dd_sum={}
    for k in set(d1.keys()).union(d2.keys()):
        dd_sum[k]=0
        if k in d1:
            dd_sum[k]+=d1[k]
        if k in d2:
            dd_sum[k]+=d2[k]
        if (dd_sum[k]==0): del dd_sum[k]
    return(dd_sum)

def dict_of_keyword_counts_in_series(series,keywords):
    series = replace_strings(series.str.lower().str.replace('[\n&,.-]',' '),
                             replace_text
                            )
    return(count_strings(series,keywords))

def dataframe_for_month(df, year=2017, month=1):
    """Yields a subset the dataframe with only those rows in the given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    if (mm==12): mm2,yyyy2=1,yyyy+1
    else: mm2,yyyy2=mm+1,yyyy

    print("Selecting rows for %i-%i" % (yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy,mm,1)) & (df.date<date(yyyy2,mm2,1)) ]
    print("Found",len(month_df.index),"rows for this month")

    return month_df


def keywords_in_dataframe(month_df,year=2017,month=1):
    """ Count each keyword in the dataframe's important columns for a given year & month
        Returns a sorted dictionary of (keyword:count)
    """
    import operator   #for sorted(), used in sorting items into OrderedDict

    print("Counting occurrences of keywords for this month")
    d_wtlma = dict_of_keyword_counts_in_series(month_df['Want to Learn More About'],vocab)
    d_actions=dict_of_keyword_counts_in_series(month_df['Action Items'],vocab)

    print("Adding together the results")
    dict_of_words_to_freq = sum_of_dicts(d_wtlma,d_actions)

    sorted_dict = collections.OrderedDict(sorted(dict_of_words_to_freq.items(), key=lambda t: t[1],reverse=True))
    return sorted_dict


def generate_wordcloud_from_dict(dict_of_words_to_freq):
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    # (to install, run "pip install wordcloud")
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt

    print("Generating the wordcloud")
    wordcloud= WordCloud(width=2500,height=500,
                         prefer_horizontal=0.9,
                         relative_scaling=0.5,
                         max_font_size=144,
                         background_color="white"
                        ).generate_from_frequencies(dict_of_words_to_freq)
    return wordcloud

def file_wordcloud_for_month(keywords_for_month, year=2017, month=8):
    wc_for_month = generate_wordcloud_from_dict(keywords_for_month)

    # Display the generated image
    plt.imshow(wc_for_month,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    plt.show()

    filename = "wordcloud-"+calendar.month_name[month]+".png"
    plt.imsave(filename,wc_for_month,format="png")

    return filename


def file_graph_for_month_kwd(kwd,kwd_pos,vals,months,line_color):
    #Put the title in the graph
    plt.title(kwd,size=24)
    #Plot the line
    plt.axis('off')
    plt.plot(vals,color=line_color,linewidth=3)
    #write the text for the labels at each end of the line
    m_minus_2_percent = "{0:.0f}%".format(vals[0] * 100)
    m_this_percent = "{0:.0f}%".format(vals[2] * 100)
    plt.annotate(calendar.month_name[months[0]]+"\n"+m_minus_2_percent,xy=(0.1,vals[0]),size=16)
    plt.annotate(calendar.month_name[months[2]]+"\n"+m_this_percent,xy=(2,vals[2]),size=16)
    #show the graph, then file it
    plt.show()
    filename = "graph-"+str(kwd_pos)+".png"
    #plt.imsave(filename,wc_for_month,format="png")
    return

def new_run_in_slide(para,text='text',fontname='Arial',fontsize=24):
    "Add new run to a paragraph in a powerpoint slide"
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.name = fontname
    return run

replace_text={
    "/": " ",                       "\n": " ",                             "big data": "big_data",
    "new stack": "new_stack",       "flexible capacity": "flex_capacity",  "gen 10": "gen10",
    "gen z": "gen-z",                "cloud cruiser": "cloudcruiser",      "store once": "storeonce",
    "smart city": "smart_city",     "future city": "smart_city",           "azure stack": "azurestack",
    "intelligent edge": "edgeline", "clear pass": "clearpass",             "one view": "oneview",
    "the machine": "the_machine",   "open stack": "openstack",             "store once": "storeonce",
    "office 365": "office365",      "hp financial services":"hpefs",       "mobility": "wireless",
    "hpe financial services": "hpefs",     "hpfs": "hpefs",                "integrity": "superdome"

}

#use this specific set of words as the vocab to look for
vocab=["iot","edgeline","smart_city", "big_data","sap","apollo", "saas", "analytics", "sgi",
       "networking", "wireless", "aruba", "arista", "clearpass", "naas", "skype", "meridian",
       "airwave","composable", "niara", "flex_capacity", "helion", "azurestack","synergy",
       "moonshot","converged","hyperconverged","hybrid","docker","vdi","new_stack","cloudline",
       "cloudsystem","easyconnect","openstack","devops","3par","bura","simplivity","nimble",
       "scality","storeonce", "oneview", "office365","pointnext","hpefs","the_machine","photonics",
       "blockchain","pathfinder","gen10","gen-z","cloudcruiser","blades","superdome"
      ]

print("Importing excel")
all_df = pd.read_excel(open('Insights.xlsx','rb'),header=8)

print("Adding structured date column")
all_df.insert(loc=0,column='date',value=all_df['Visit Date'].apply(make_date))

## Given this month and year, calc the number of the previous two months
mm,yyyy=8,2017

if (mm==1): mm_minus_1,year_for_mm_minus_1 = 12, yyyy-1
else: mm_minus_1, year_for_mm_minus_1 = mm-1, yyyy

if (mm<=2): mm_minus_2,year_for_mm_minus_2 = mm+10, yyyy-1
else: mm_minus_2, year_for_mm_minus_2 = mm-2, yyyy

## Now do the counting & work: count the keywords for this month and the previous two,
## and build their wordclouds.
df_for_month = dataframe_for_month(all_df, year=yyyy, month=mm)
kwd_count_for_month = keywords_in_dataframe(df_for_month,year=yyyy,month=mm)
print("Top keywords and counts found for month",mm,":",
      list(kwd_count_for_month.items())[:5] )
file_wordcloud_for_month(kwd_count_for_month,
                         year=yyyy,month=mm)

df_for_month_minus_1 = dataframe_for_month(all_df,
                                           year=year_for_mm_minus_1, month=mm_minus_1)
kwd_count_for_m_minus_1 = keywords_in_dataframe(df_for_month_minus_1,
                                                   year=year_for_mm_minus_1,month=mm_minus_1)
print("Top keywords and counts found for month",mm_minus_1,":",
      list(kwd_count_for_m_minus_1.items())[:5] )
file_wordcloud_for_month(kwd_count_for_m_minus_1,
                         year=year_for_mm_minus_1,month=mm_minus_1)

df_for_month_minus_2 = dataframe_for_month(all_df,
                                           year=year_for_mm_minus_2, month=mm_minus_2)
kwd_count_for_m_minus_2 = keywords_in_dataframe(df_for_month_minus_2,
                                                   year=year_for_mm_minus_2,month=mm_minus_2)
print("Top keywords and counts found for month",mm_minus_2,":",
      list(kwd_count_for_m_minus_2.items())[:5] )
file_wordcloud_for_month(kwd_count_for_m_minus_2,
                         year=year_for_mm_minus_2,month=mm_minus_2)


## Get top 3 keywords for this month, and graph their usage
kwd0,kwd1,kwd2 = list(kwd_count_for_month)[:3]
print("Top 3 keywords for current month are",kwd0,kwd1,kwd2)

import matplotlib.pyplot as plt
months=[mm_minus_2,mm_minus_1,mm]
vals=[list(kwd_count_for_m_minus_2.values())[0]/len(df_for_month_minus_2.index),
      list(kwd_count_for_m_minus_1.values())[0]/len(df_for_month_minus_1.index),
      list(kwd_count_for_month.values())[0]/len(df_for_month.index)]
file_graph_for_month_kwd(kwd0,"1st",vals,months,'#009000')

vals=[list(kwd_count_for_m_minus_2.values())[1]/len(df_for_month_minus_2.index),
      list(kwd_count_for_m_minus_1.values())[1]/len(df_for_month_minus_1.index),
      list(kwd_count_for_month.values())[1]/len(df_for_month.index)]
file_graph_for_month_kwd(kwd1,"2nd",vals,months,'#000090')

vals=[list(kwd_count_for_m_minus_2.values())[2]/len(df_for_month_minus_2.index),
      list(kwd_count_for_m_minus_1.values())[2]/len(df_for_month_minus_1.index),
      list(kwd_count_for_month.values())[2]/len(df_for_month.index)]
file_graph_for_month_kwd(kwd2,"3rd",vals,months,'#900000')

### Now generate the powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt

# These appear to be the layouts used in the master for this slide deck
LAYOUT_TITLE_WITH_DARK_PICTURE = 0
LAYOUT_TITLE_SLIDE_WITH_NAME   = 1
LAYOUT_DIVIDER                 = 2
LAYOUT_TITLE_ONLY              = 3
LAYOUT_TITLE_AND_SUBTITLE      = 4
LAYOUT_BLANK                   = 5

## Open up the source presentation
prs = Presentation('GCA_Customer_Insights_Month-Year.pptx')
this_month = calendar.month_name[mm]
print("Creating presentation for",this_month)

## Modify existing title slide with month
print("Modifying text in slide 0")
s = prs.slides[0]
slide_shapes = s.shapes
text_frame = slide_shapes[0].text_frame  # should be the the title textframe
#clear existing text, and write new text into title textframe
text_frame.clear()
# First para with a run of 60-point Arial font
new_run_in_slide(text_frame.paragraphs[0],text='Customer Insights',fontname='Arial',fontsize=60)
# Second para with a run of 32-point font
new_run_in_slide(text_frame.add_paragraph(),text='Learnings from '+this_month+' EBC/CEC visits',fontsize=32)

## Modifying main wordcloud slide by changing title and adding pic for this month's wordcloud
print("Modifying text and adding wordcloud in slide 2")
s = prs.slides[2]
slide_shapes=s.shapes
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="In "+this_month+" customers wanted to learn more about...",
       fontname="Arial",fontsize=28)
left=Inches(0.5); top=Inches(4.0)
slide_shapes.add_picture("wordcloud-"+this_month+".png",left,top,height=Inches(2.5))


## Modifying 3-month wordcloud slide by adding pic for last 3 months' wordcloud
print("Adding three wordclouds in slide 3")
s = prs.slides[3]
slide_shapes=s.shapes
left=Inches(3.65)
top=Inches(1.4)
slide_shapes.add_picture("wordcloud-"+calendar.month_name[mm_minus_2]+".png",left,top,height=Inches(1.5))
top=Inches(3.25)
slide_shapes.add_picture("wordcloud-"+calendar.month_name[mm_minus_1]+".png",
                         left,top,height=Inches(1.5))
top=Inches(5.05)
slide_shapes.add_picture("wordcloud-"+this_month+".png",
                         left,top,height=Inches(1.5))

## Close the source presentation
print("Saving Powerpoint file for ",this_month)
prs.save('GCA_Customer_Insights_'+this_month+'-'+str(yyyy)+'.pptx')
