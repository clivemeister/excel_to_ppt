import pandas as pd
from datetime import datetime, date
import matplotlib.pyplot as plt
import collections
from collections import Counter
import calendar
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import os
import sys
import logging
import re
import sys, getopt

logger = logging.getLogger(__name__)
##logger.setLevel(logging.WARNING)
console=logging.StreamHandler()
console.setLevel(logging.WARNING)
formatter=logging.Formatter('%(asctime)s %(levelname)s %(message)s')
console.setFormatter(formatter)
logger.addHandler(console)

tmpdir = "tmp/"
icondir = "icons/"
excel_file='Insights.xlsx'
stop_after_wordcheck = False
yyyy,mm=date.today().year,date.today().month
which_ctr='PA'      #default

def print_help():
    print("single_centre_view.py  --ifile=<inputExcelFile>    default is Insights.xlsx")
    print("                       --ctr=<cc>                 centre code: PA, H, LON1, NY1, SNG.  Default is PA.")
    print("                       --year=<yyyy>              year to process, default is current year")
    print("                       --month=<mm>               month to process, default is current month")
    print("                       -d                     turn on debugging trace")
    print("                       -i                     turn on information trace")
    print("e.g   single_centre_view -i --ctr=SNG --month=2")
    return

if __name__=="__main__":
    logging.debug("Parsing arguments")
    try:
        opts, args = getopt.getopt(sys.argv[1:],"hdi",["ifile=","year=","month=","ctr="])
    except getopt.GetoptError:
        print_help()
        sys.exit(2)
    for opt, arg in opts:
        if opt == "-h":
            print_help()
            sys.exit()
        elif opt == "-d":
            ##console.setLevel(logging.DEBUG)
            logger.setLevel(logging.DEBUG)
        elif opt == "-i":
            ##console.setLevel(logging.INFO)
            logger.setLevel(logging.INFO)
        elif opt in ("--ifile"):
            excel_file = arg
        elif opt in ("-y:","--year"):
            logging.debug("Found argument yyyy with {}".format(arg))
            yyyy = int(arg)
        elif opt in ("-m:", "--month"):
            logging.debug("Found argument mm with {}".format(arg))
            mm = int(arg)
        elif opt in ("-c:", "--ctr"):
            logging.debug("Found argument ctr with {}".format(arg))
            which_ctr = arg
        else:
            assert False, "unhandled option"


logger.info('Starting run for %i-%i for %s...' % (yyyy,mm,which_ctr))

# Read in the ini file
import configparser
import json
ini_file = 'excel_to_ppt.ini'
cfg = configparser.ConfigParser()
cfg.optionxform = str    # read strings as-is from INI file (default is to lowercase keys)
cfg.read(ini_file)

colour_list=[]   # list of the colour code values
if 'colours' in cfg:
    colour_codes = dict(cfg.items('colours'))
    for i in cfg.items('colours'):
        colour_list.append(i[1])       # put this hex value in list of available colours
else:
    logger.error('No [colours] section in {}'.format(ini_file))

if 'keywords' in cfg:
    c_to_k=cfg.items('keywords')
    dict_colour_of_keywords={}
    for i in cfg.items('keywords'):
        if i[0] in colour_codes:
            dict_colour_of_keywords[colour_codes[i[0]]]=json.loads(i[1])
        else:
            logger.error('In {}, section [keywords], found colour {} which was not listed in [colours] section'.format(ini_file,i[0]))
else:
    logger.error('No [keywords] section in {}'.format(ini_file))

#Build a list of all the keywords.  Use all the things we have a colour for as the list items.
vocab=[]
for i in dict_colour_of_keywords.values():
    for j in i:
        vocab.append(j)

# load the list of synonyms (or mis-spellings) of the keywords
synonym_list={}
for i in cfg.items('synonyms'):
    lst=json.loads(i[1])
    for j in lst:
        synonym_list[j]=i[0]

JapanAndChinaToOther=True
industry_list=[]   # list of the industry code values
if 'industries' in cfg:
    industry_longnames = dict(cfg.items('industries'))
    for i in cfg.items('industries'):
        industry_list.append(i[0])
    if (cfg.has_option('industries','JapanAndChinaToOther')):
        try:
            JapanAndChinaToOther = cfg.getboolean('industries','JapanAndChinaToOther')
        except:
            logger.error('In [industries] section item JapanAndChinaToOther is not boolean')
else:
    logger.error('No [industries] section in {}'.format(ini_file))

centres = ["H","NY1","SNG","LON1","PA"]   # short names used in the Excel file
centres_long = ["Houston", "New York", "Singapore", "London", "Palo Alto"]  # same order as short names

stop_words=[]
if 'stopwords' in cfg:
    stop_words+=json.loads(cfg.get('stopwords','stop_words'))
else:
    logger.error('No [stopwords] section in {}'.format(ini_file))
stop_words+=vocab



def make_date(cell_val):
    """Used to create a better-structured date value from the sometimes-odd values in the date column
       Returns a datetime
    """
    if type(cell_val) is datetime:
        v=cell_val.date()
    elif type(cell_val) is str:
        try:
            v=datetime.strptime(cell_val,"%b %d, %Y").date()
        except ValueError:
            v=date.fromordinal(1)
    else:
        v=date.fromordinal(2)
    return v

def tidy_text(cell_val):
    """Standardises the text in a cell:
       Returns the tidied text
    """
    if type(cell_val) is str:
        cell = cell_val.lower()
        cell=re.sub('[\n&@,.:-]',' ',cell)
        cell=" ".join(cell.split())   # idiom to turn multiple spaces between words into single spaces
        for k,v in synonym_list.items():
            cell=cell.replace(k,v)
    else:
        cell=""
    return cell

def replace_strings(series,repl_dict):
    for k,v in repl_dict.items():
        series=series.str.replace(k,v)
    return(series)

def bool_list_of_occurrences(series,kwd):
    """Return a boolean list, 1 where an element of 'series' contains word 'kwd', else 0
       Returns a list of 1s and 0s
    """
    pattern = r'\b{0}\b'.format(kwd)
    bool_list=series.str.contains(pattern)
    return bool_list

def keyword_counts_in_series(series,keywords):
    """Return a Counter counting how often each word in 'keywords' appears in 'series'
       Actually we count the number of cells in 'series' with keyword in it, so multiple
       occurences per cell only count once.
       Returns a Counter.
    """
    if keywords is str:
        #get here only if keywords is a single string, not a list
        count_list=sum(bool_list_of_occurrences(series,keywords))
    else:
        count_list=Counter()
        for kwd in keywords:
            count_list[kwd]=sum(bool_list_of_occurrences(series,keywords))
    return count_list

def dataframe_for_month(df, year=2018, month=1):
    """Yields a subset the dataframe with only those rows in the given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    if (mm==12): mm2,yyyy2=1,yyyy+1
    else: mm2,yyyy2=mm+1,yyyy

    logger.debug("Selecting rows for %i-%i" % (yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy,mm,1)) & (df.date<date(yyyy2,mm2,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def dataframe_for_6months(df, year=2018, month=1):
    """Yields a subset the dataframe with those rows for last 6 months up to given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    # calc month after this one, to set upper limit for search
    if (mm==12): mm_end,yyyy_end=1,yyyy+1
    else: mm_end,yyyy_end=mm+1,yyyy
    # calc 6 months ago, to set lower limit for search
    if (mm<6): mm_start,yyyy_start=mm+7,yyyy-1
    else: mm_start,yyyy_start=mm-5,yyyy

    logger.debug("Selecting rows for %i-%i to %i-%i" % (yyyy_start,mm_start,yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy_start,mm_start,1)) & (df.date<date(yyyy_end,mm_end,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def found_word_list(df, kwd):
    """Passed a dataframe with columns of at least 'wtlma' and 'ai', together with a keyword
       Return a list (actually a Series) with True wherever the dataframe contains kwd or a synonym in the relevant columns,
       and False otherwise.  Can be used to subset the original dataframe to pick out the rows with the keyword in them:
           df[found_word_list(df,"foo")]   => subset of df with "foo" in one of the columns
    """
    strings_wtlma = replace_strings(df.wtlma, synonym_list)
    strings_ai = replace_strings(df.ai, synonym_list)
    found_list = bool_list_of_occurrences(strings_wtlma,kwd) | bool_list_of_occurrences(strings_ai,kwd)
    return found_list

def keywords_in_dataframe(df,keyword_list):
    """ Count each keyword in the dataframe's important columns
        Returns a Counter collection of (keyword:count)
    """
    import operator   #for sorted(), used in sorting items into OrderedDict
    counter_words_to_freq = Counter()
    for w in keyword_list:
         counter_words_to_freq[w] += sum(found_word_list(df,w))

    return counter_words_to_freq

def count_rows_with_comments(df):
    """ Count the number of rows in dataframe with a comment in either <Want to Learn More About> or <Action Items>
        Returns the count
    """
    return (df["Want to Learn More About"].notnull() | df["Action Items"].notnull()).sum()

def counts_by_centre(dataframe):
    centre_counts = dataframe.Ctr.value_counts()
    counts_list = (centre_counts.PA if hasattr(centre_counts,"PA") else 0,
                   centre_counts.H if hasattr(centre_counts,"H") else 0,
                   centre_counts.NY1 if hasattr(centre_counts,"NY1") else 0,
                   centre_counts.LON1 if hasattr(centre_counts,"LON1") else 0,
                   centre_counts.SNG if hasattr(centre_counts,"SNG") else 0
                   )
    return counts_list

class SimpleGroupedColorFunc(object):
    """Create a color function object which assigns EXACT colors
       to certain words based on the color to words mapping

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        self.word_to_color = {word: color
                              for (color, words) in color_to_words.items()
                              for word in words}

        self.default_color = default_color

    def __call__(self, word, **kwargs):
        return self.word_to_color.get(word, self.default_color)

class GroupedColorFunc(object):
    """Create a color function object which assigns DIFFERENT SHADES of
       specified colors to certain words based on the color to words mapping.

       Uses wordcloud.get_single_color_func

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        from wordcloud import get_single_color_func
        self.color_func_to_words = [
            (get_single_color_func(color), set(words))
            for (color, words) in color_to_words.items()]

        self.default_color_func = get_single_color_func(default_color)

    def get_color_func(self, word):
        """Returns a single_color_func associated with the word"""
        try:
            color_func = next(
                color_func for (color_func, words) in self.color_func_to_words
                if word in words)
        except StopIteration:
            color_func = self.default_color_func

        return color_func

    def __call__(self, word, **kwargs):
        return self.get_color_func(word)(word, **kwargs)

def file_wordcloud_for_month(keywords_for_month, useful_rows_for_month, year=2017, month=8):
    """Input is a keywords_for_month, a Counter, plus useful_rows_for_month, an integer, and year and month as integers
       Returns the filename where the wordcloud is stored
    """
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    # (to install, run "pip install wordcloud"
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt
    import os

    # set font so 30% occurrence of top word uses 196 point font
    percent = round(100*(keywords_for_month.most_common(1)[0][1]/useful_rows_for_month))
    font_for_biggest_word = round( 196 * percent/30 )
    logger.debug("font for word {} is {} based on {}% = {} / {}".format(
                               keywords_for_month.most_common(1)[0][0],
                                     font_for_biggest_word,
                                             percent,
                                                keywords_for_month.most_common(1)[0][1],
                                                        useful_rows_for_month))

    font_path = os.path.join("C:",os.sep,"Windows",os.sep,"Fonts",os.sep,'arial.ttf')
    logger.debug("Generating the wordcloud")
    keywords_for_month += Counter()    # remove any zero or negative counts from the list
    wc_for_month = WordCloud(font_path=font_path,
                             width=2500,height=500,
                             prefer_horizontal=1.0,
                             relative_scaling=0.7,
                             max_font_size=font_for_biggest_word,
                             background_color="white",
                             random_state=1
                             ).generate_from_frequencies(keywords_for_month)


    # Words that are not in any of the dict_colour_of_keywords values
    # will be colored with a grey single color function
    default_color = 'grey'

    # Apply our color function
    #grouped_color_func = GroupedColorFunc(dict_colour_of_keywords, default_color)
    grouped_color_func = SimpleGroupedColorFunc(dict_colour_of_keywords, default_color)
    wc_for_month.recolor(color_func=grouped_color_func)

    # Build the generated image
    plt.imshow(wc_for_month,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    #plt.show()

    filename = tmpdir+"wordcloud-"+calendar.month_name[month]+".png"
    plt.imsave(filename,wc_for_month,format="png")
    plt.close()

    return filename


def file_graph_for_month_kwd(kwd,kwd_pos,vals,months,line_color):
    fig,ax=plt.subplots(figsize=(5.75,3.25))
    ##############ax.clear()

    ax.plot(vals,line_color,linewidth=3)
    ax.set_axis_off()
    plt.ylim(min(vals)-0.05,max(vals)+0.05)
    m_minus_2_percent = "{0:.0f}%".format(vals[0] * 100)
    m_this_percent    = "{0:.0f}%".format(vals[2] * 100)
    ax.text(0,vals[0]+0.02,calendar.month_abbr[months[0]]+"\n"+m_minus_2_percent,fontsize=30)
    ax.text(2,vals[2]+0.02,calendar.month_abbr[months[2]]+"\n"+m_this_percent,fontsize=30)
    filename = tmpdir+"graph-"+str(kwd_pos)+".png"
    logger.debug("Saving %s graph for keyword %s in file %s" % (kwd_pos,kwd,filename))
    fig.savefig(filename,bbox_inches="tight")
    plt.close(fig)

    return filename

def file_donut_pie_for_month(values,name):
    fig, ax = plt.subplots()
    ax.axis('equal')
    outside, _ = ax.pie(values,startangle=90,counterclock=False,
                        colors=list(colour_list))
    ax.legend(("Palo Alto","Houston","NY","London","Singapore"),fontsize=24,bbox_to_anchor=(0.8,1.0),frameon=False)
    width = 0.50  #determines the thickness of donut rim
    plt.setp( outside, width=width, edgecolor='white')

    filename = tmpdir+"donut-"+re.sub(r"[& ]","_",str(name))+".png"
    logger.debug("Saving <%s> donut in file %s with values %r" % (name, filename, values))
    fig.savefig(filename,bbox_inches="tight")

    plt.close(fig)

    return filename

def file_donut_pie_for_industries(industries,center=""):

    if (JapanAndChinaToOther and 'Japan' in industries):
        industries.Other += industries.Japan
        del industries['Japan']

    if (JapanAndChinaToOther and 'China' in industries):
        industries.Other += industries.China
        del industries['China']

    fig, ax = plt.subplots()
    ax.axis('equal')
    width = 2.0
    outside, _ = ax.pie(industries.values,
                        startangle=90,counterclock=False,
                        colors=list(colour_list),
                        radius=5.0)
    ax.legend(industries.index.tolist(),ncol=3,fontsize=24,loc=10,bbox_to_anchor=(0.5,-2),frameon=False)
    plt.setp( outside, width=width, edgecolor='white')

    filename = tmpdir+"donut-industries"+center+".png"
    logger.debug("Saving Industries donut in file %s with values %r" % (filename, industries))
    fig.savefig(filename,bbox_inches="tight")

    plt.close(fig)

    return filename

def write_customer_list(df_for_kwd,text_frame):
    text_frame.text = "Briefings"
    font=text_frame.paragraphs[0].font
    font.size = Pt(12)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for i in (df_for_kwd)["Account Name"].str.title().iteritems():
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        new_para.text=i[1]

    return

def write_top_keywords(text_frame,titleText,kwd_counts_for_ind,rowcount,percent=True,cutoff=0):
    """Write a section headed by titleText, with bullets of the top 4 elements of Counter
       kwd_counts_for_ind plus their percentage of use (if percent=True), where
       the percentage is calculated with rowcount as the denominator.  Stop before top
       4 elements if they start to fall below the cutoff percentage.
         e.g:
           c=Counter([("alpha":5),("bravo":4),("charlie":3),("delta":1))])
           write_top_keywords(tf,"My list",c,10,percent=True,cutoff=20)
         yields:
           My list
            - alpha - 50%
            - bravo - 40%
            - charlie - 30%
    """
    text_frame.text = titleText
    font=text_frame.paragraphs[0].font
    font.size = Pt(10)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for k,v in kwd_counts_for_ind.most_common(4):
        if (100*v/rowcount>=cutoff):
            new_para = text_frame.add_paragraph()
            font=new_para.font
            font.size = Pt(8)
            font.color.theme_color = MSO_THEME_COLOR.DARK_2
            if (percent==True):
                new_para.text=" - "+k+" - "+"{0:.0f}%".format(v/rowcount * 100)
            else:
                new_para.text=" - "+k

    return

def find_text_in_shapes(slide_shapes,searchword):
    """Find the given text in the list of powerpoint shapes passed in slide_shapes
    """
    found_idx = -1
    for i in range(len(slide_shapes)):
        if (slide_shapes[i].has_text_frame):
            text_frame = slide_shapes[i].text_frame
            if (text_frame.text.find(searchword)>=0):
                found_idx = i
                break

    return found_idx

def replace_text_in_shape(slide_shapes,find,use,slidename):
    """Find the placeholder text "find" in the slide_shapes, replace it with "use",
       reporting an error on "slidename" if we can't find "find"
    """
    i = find_text_in_shapes(slide_shapes,find)
    if (i>=0):
        text_frame = slide_shapes[i].text_frame
        text_frame.text = use
    else:
        logger.error("Could not find %s placeholder on %s" % (find,slidename))

    return

def new_run_in_slide(para,text='text',fontname='Arial',fontsize=24):
    "Add new run to a paragraph in a powerpoint slide"
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.name = fontname
    return run

import struct
import imghdr

def get_image_size(fname):
    '''Determine the image type of fname and return its width, height.
    from draco (via stackoverflow)'''
    with open(fname, 'rb') as fhandle:
        head = fhandle.read(24)
        if len(head) != 24:
            return
        if imghdr.what(fname) == 'png':
            check = struct.unpack('>i', head[4:8])[0]
            if check != 0x0d0a1a0a:
                return
            width, height = struct.unpack('>ii', head[16:24])
        elif imghdr.what(fname) == 'gif':
            width, height = struct.unpack('<HH', head[6:10])
        elif imghdr.what(fname) == 'jpeg':
            try:
                fhandle.seek(0) # Read 0xff next
                size = 2
                ftype = 0
                while not 0xc0 <= ftype <= 0xcf:
                    fhandle.seek(size, 1)
                    byte = fhandle.read(1)
                    while ord(byte) == 0xff:
                        byte = fhandle.read(1)
                    ftype = ord(byte)
                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
                # We are at a SOFn block
                fhandle.seek(1, 1)  # Skip `precision' byte.
                height, width = struct.unpack('>HH', fhandle.read(4))
            except Exception: #IGNORE:W0703
                return
        else:
            return
        return width, height

def add_icon(ss,icon_name,left,top,small=False):
    ## ok if icon exists, and use error icon if not
    logger.debug("adding icon {}".format(icon_name))
    try:
        filename = icondir+icon_name+".png"
        icon_width, icon_height = get_image_size(filename)
        pixels_per_mm = 6
        if small:
            icon_width = 0.4*icon_width
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(12),Mm(12))
        else:
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(34),Mm(34))
    except:
        logger.warning("Could not find icon {}".format(icondir+icon_name+'.png'))
    return

logger.info("Importing excel file "+excel_file)
all_df = pd.read_excel(open(excel_file,'rb'),header=8,usecols="A:S")

logger.debug("Adding structured columns")
all_df.insert(loc=0,column='date',value=all_df['Visit Date'].apply(make_date))
all_df.insert(loc=1,column='wtlma',value=all_df['Want to Learn More About'].apply(tidy_text))
all_df.insert(loc=2,column='ai',value=all_df['Action Items'].apply(tidy_text))


## Given the month and year, calc the number of the previous few months
if (mm==1): mm_minus_1,year_for_mm_minus_1 = 12, yyyy-1
else: mm_minus_1, year_for_mm_minus_1 = mm-1, yyyy
if (mm<=2): mm_minus_2,year_for_mm_minus_2 = mm+10, yyyy-1
else: mm_minus_2, year_for_mm_minus_2 = mm-2, yyyy
if (mm<6): mm_6_before,year_for_mm_6_before = mm+7, yyyy-1   #6 months before Jan is Aug
else: mm_6_before, year_for_mm_6_before = mm-5, yyyy    #6 months before Jul is Feb

### Now start to generate the powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt, Mm

# These appear to be the layouts used in the master for this slide deck
LAYOUT_TITLE_WITH_DARK_PICTURE = 0
LAYOUT_TITLE_SLIDE_WITH_NAME   = 1
LAYOUT_DIVIDER                 = 2
LAYOUT_TITLE_ONLY              = 3
LAYOUT_TITLE_AND_SUBTITLE      = 4
LAYOUT_BLANK                   = 5

## Open up the source presentation
prs = Presentation('GCA_Centre_Insights_Month-Year - ctr.pptx')
this_month = calendar.month_name[mm]
earliest_month = calendar.month_name[mm_minus_2]
logger.info("Creating presentation for %s" % (this_month))


################################################
## Build Top interests and industries for last 6 months
################################################
logger.info(">>>> Top 5 interests, top 3 industries, and their top interests, by centre, for last 6 months")
df_6months = dataframe_for_6months(all_df, year=yyyy, month=mm)

#build a dictionary of dataframes subsetted by centre, a dictionary of keywords by centre,
#and a dict of top 3 industries per centre and their keywords (where the key is a tuple of (centre, industry) )
dfs_6m_ctr={}
kwd_counts_6m_ctr={}
industry_counts_6m = df_6months["Industry"].value_counts()
kwd_counts_6m_for_top_inds={}
commented_rows_for_6m={}

for c in centres:
    logger.info("Working on counts for {}".format(c))
    dfs_6m_ctr[c]=df_6months[df_6months.Ctr==c]
    #First, the top keywords for that Centre
    kwd_counts_6m_ctr[c]=keywords_in_dataframe(dfs_6m_ctr[c],vocab)
    logger.debug("Top keyword/counts %s: %r" % (c,kwd_counts_6m_ctr[c].most_common(5)) )
    #Now the top keywords in each industry for that centre
    industry_counts_6m[c] = (dfs_6m_ctr[c])["Industry"].value_counts()
    for ind in (industry_counts_6m[c]).index:
        this_df = dfs_6m_ctr[c][ dfs_6m_ctr[c]["Industry"]==ind ]
        kwd_counts_6m_for_top_inds[(c,ind)]=keywords_in_dataframe(this_df, vocab)
        commented_rows_for_6m[(c,ind)]=count_rows_with_comments(this_df)
        logger.debug("Top keyword/counts in %s for %s: %r" % (c,ind,kwd_counts_6m_for_top_inds[(c,ind)].most_common(3)) )


################################################
## The slide: EBC specific volumes & interests for last 6 months
################################################
logger.info(">>>> slide: %s specific industry volumes and top interests" % (which_ctr))
## Generate the image for the big donut showing volume for all industries across the months
ctr_industry_counts = dfs_6m_ctr[which_ctr]["Industry"].value_counts()
file_donut_ctr_ind_vols = file_donut_pie_for_industries(ctr_industry_counts,center=which_ctr)

## Build the slide
logger.info("Setting the title ")
s = prs.slides[0]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text=which_ctr+" six month view ("+calendar.month_name[mm_6_before]+"-"+this_month+")",
       fontname="Arial",fontsize=28)

#Add the one big donut showing volumes for each industry to the slide
logger.info("Adding main donut to six month view slide")
left=Mm(43); top=Mm(48)
slide_shapes.add_picture(file_donut_ctr_ind_vols,left,top,height=Mm(115),width=Mm(94))

#Update, by centre, the top interests, and the top industries with their top interests
logger.info("Setting the per-centre top 5 interests, together with per-centre top 3 industries and their interests")
left_pos=[Mm(180),Mm(210),Mm(240),Mm(270),Mm(299)]   #distances to centre of icon for each row
#First, the top interests for PA
for col,p in enumerate(kwd_counts_6m_ctr[which_ctr].most_common(5)):
    # p is (keyword: count) for each of the keywords, so p[0] is the keyword itself
    replace_text_in_shape(slide_shapes,"interest-{}".format(col),p[0],"Centre slide")
    add_icon(slide_shapes,p[0],left=left_pos[col],top=Mm(59),small=True)
#Next, the top industries with their interests for this centre - industry_counts_6m[c] is already ordered highest->lowest count
n=0  #count how many displayed - need to do this separately from the loop count, as we ignore "Other" as an industry group
for ind in industry_counts_6m[which_ctr].index:
    if (ind!="Other"):
        logger.debug("For this centre, industry <%i> is <%s>" %(n,ind))
        #Write the list of top interests for this industry
        idx = find_text_in_shapes(slide_shapes,"industry-{}".format(n))
        if (idx>=0):
            write_top_keywords(slide_shapes[idx].text_frame,
                               ind,
                               kwd_counts_6m_for_top_inds[(which_ctr,ind)],
                               commented_rows_for_6m[(which_ctr,ind)]
                              )
        else:
            logger.error("Could not find <industry-{}> placeholder on slide - idx={}".format(n,idx))
        n+=1   #increment the number of industries written out for this centre
    if (n>=3): break   #exit the loop after writing in 3 industries+interests


################################################
## Close the source presentation
################################################
logger.info("Saving Powerpoint file for "+this_month)
prs.save('GCA_Centre_Insights_'+this_month+'-'+str(yyyy)+"-"+which_ctr+'.pptx')
## Close any open figures
plt.close("all")
logger.info("...and we're done!")
for h in list(logger.handlers): logger.removeHandler(h)   # may be several here if we've crashed sometimes
