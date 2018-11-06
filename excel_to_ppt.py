""" Read the insights.xls file and conver to GCA_Customer_Insights_mmmm-yyyy.pptx

    Typical usage to get Insights for June 2018 would be:
        py excel_to_ppt.py -m6 -y2018 

    Uses the GCA_Customer_Insights_Month-Year.pptx as a template.
    Uses the excel_to_ppt.ini file for the vocabulary of words, synonyms, colours to use, etc etc

    Required packages: pandas, matplotlib, python-pptx, xlrd, wordcloud, vaderSentiment, requests
    You can use 'pip install' to get these down.
"""
import pandas as pd
from datetime import datetime, date
import matplotlib.pyplot as plt
import collections
from collections import Counter
import calendar
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import os
import sys
import logging
import re
import sys, getopt

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

tmpdir = "tmp/"
icondir = "icons/"
excel_file='Insights.xlsx'
stop_after_wordcheck = False
yyyy,mm=date.today().year,date.today().month-1
header_row=1
oldMode_1 = False
sentimentCalcs = False

def print_help():
    print("excel_to_ppt.py  -ifile=<inputExcelFile>    default is Insights.xlsx")
    print("                 --year=<yyyy>              year to process, default is current year")
    print("                 --month=<mm>               month to process, default is current month")
    print("                 -m<mm>                     month to process, e.g. -m7")
    print("                 -y<yyyy>                   year to process, e.g. -y2018")
    print("                 -w                     stop after showing possible extra words")
    print("                 -r<n>                  set which row of the spreadsheet is the first (header) row.  Default is 1.")
    print("                 -d                     turn on debugging trace")
    print("                 -s                     turn on experiemental sentiment analysis")
    print("                 -q                     turn on quiet mode - shows only information")
    print("For example, excel_to_ppt.py -m8 -y2018 -s")
    return

if __name__=="__main__":
    logging.debug("Parsing arguments")
    try:
        opts, args = getopt.getopt(sys.argv[1:],"?hdwsvi:y:m:r:",["ifile=","year=","month="])
    except getopt.GetoptError as err:
        print(err)
        print_help()
        sys.exit(2)
    logger.debug("opt: "+str(opts)+" and args: "+str(args))  
    for opt, arg in opts:
        if opt in ("-?","-h"):
            print_help()
            sys.exit()
        elif opt == "-d":
            logger.setLevel(logging.DEBUG)
        elif opt == "-q": # lowest level of messages printed
            logger.setLevel(logging.WARNING)
        elif opt in ("--ifile"):
            excel_file = arg
        elif opt in ("-y","--year"):
            logging.debug("Found argument yyyy with {}".format(arg))
            yyyy = int(arg)
        elif opt in ("-m", "--month"):
            logging.debug("Found argument mm with {}".format(arg))
            mm = int(arg)
        elif opt in ("-r"):
            logging.debug("Found argument header_row with {}".format(arg))
            header_row = int(arg)
        elif opt == "-w":
            stop_after_wordcheck = True
        elif opt == "-s":
            sentimentCalcs = True
        elif opt == "-o":
            logging.info("Working in old mode {}".format(arg))
            if int(arg)==1:
                oldMode_1 = True
            #endif
        #endif

logger.info('Starting run for %i-%i...' % (yyyy,mm))

# Read in the ini file
import configparser
import json
ini_file = 'excel_to_ppt.ini'
cfg = configparser.ConfigParser()
cfg.optionxform = str    # read strings as-is from INI file (default is to lowercase keys)
cfg.read(ini_file)

colour_list=[]   # list of the colour code values
if 'colours' in cfg:
    colour_codes = dict(cfg.items('colours'))
    for i in cfg.items('colours'):
        colour_list.append(i[1])       # put this hex value in list of available colours
else:
    logger.error('No [colours] section in {}'.format(ini_file))

if 'keywords' in cfg:
    c_to_k=cfg.items('keywords')
    dict_colour_of_keywords={}
    for i in cfg.items('keywords'):
        if i[0] in colour_codes:
            dict_colour_of_keywords[colour_codes[i[0]]]=json.loads(i[1])
        else:
            logger.error('In {}, section [keywords], found colour {} which was not listed in [colours] section'.format(ini_file,i[0]))
else:
    logger.error('No [keywords] section in {}'.format(ini_file))

#Build a list of all the keywords.  Use all the things we have a colour for as the list items.
vocab=[]
for i in dict_colour_of_keywords.values():
    for j in i:
        vocab.append(j)

# load the list of synonyms (or mis-spellings) of the keywords
synonym_list={}
for i in cfg.items('synonyms'):
    lst=json.loads(i[1])
    for j in lst:
        synonym_list[j]=i[0]

JapanAndChinaToOther=True
industry_list=[]   # list of the industry code values
if 'industries' in cfg:
    industry_longnames = dict(cfg.items('industries'))
    for i in cfg.items('industries'):
        industry_list.append(i[0])
    if (cfg.has_option('industries','JapanAndChinaToOther')):
        try:
            JapanAndChinaToOther = cfg.getboolean('industries','JapanAndChinaToOther')
        except:
            logger.error('In [industries] section item JapanAndChinaToOther is not boolean')
else:
    logger.error('No [industries] section in {}'.format(ini_file))

centres = ["H","NY1","SNG","LON1","PA"]   # short names used in the Excel file
centres_long = ["Houston", "New York", "Singapore", "London", "Palo Alto"]  # same order as short names

stop_words=[]
if 'stopwords' in cfg:
    stop_words+=json.loads(cfg.get('stopwords','stop_words'))
else:
    logger.error('No [stopwords] section in {}'.format(ini_file))
stop_words+=vocab



def make_date(cell_val):
    """Used to create a better-structured date value from the sometimes-odd values in the date column
       Returns a datetime
    """
    if type(cell_val) is datetime:
        v=cell_val.date()
    elif type(cell_val) is str:
        try:
            v=datetime.strptime(cell_val,"%b %d, %Y").date()
        except ValueError:
            v=date.fromordinal(1)
    else:
        v=date.fromordinal(2)
    return v

def tidy_text(cell_val):
    """Standardises the text in a cell: removes lots of punctuation, and replaces synonyms by their root word
       Returns the tidied text
    """
    if type(cell_val) is str:
        cell = cell_val.lower()
        cell=re.sub('[\n&@,.:-]',' ',cell)
        cell=" ".join(cell.split())   # idiom to turn multiple spaces between words into single spaces
        for k,v in synonym_list.items():
            cell=cell.replace(k,v)
    else:
        cell=""
    return cell

def bool_list_of_occurrences(series,kwd):
    """Return a boolean list, 1 where an element of 'series' contains word 'kwd', else 0
       Returns a list of 1s and 0s
    """
    pattern = r'\b{0}\b'.format(kwd)
    bool_list=series.str.contains(pattern)
    return bool_list

def keyword_counts_in_series(series,keywords):
    """Return a Counter counting how often each word in 'keywords' appears in 'series'
       Actually we count the number of cells in 'series' with keyword in it, so multiple
       occurences per cell only count once.
       Returns a Counter.
    """
    if keywords is str:
        #get here only if keywords is a single string, not a list
        count_list=sum(bool_list_of_occurrences(series,keywords))
    else:
        count_list=Counter()
        for kwd in keywords:
            count_list[kwd]=sum(bool_list_of_occurrences(series,keywords))
    return count_list

def print_new_candidate_words(df,stop_words,top_n=25):
    """Print the top_n words in the relevant columns in the dataframe df that aren't in the
       stop_words
       Input are a dataframe with columns 'wtlma' and 'actions',
       a list in stop_words of single words to ignore, and the number of top words to show (default 25)
    """
    big_lst=[]
    for lst in list(df.wtlma.str.split()):
        if type(lst)==list: big_lst += lst
    for lst in list(df.actions.str.split()):
        if type(lst)==list: big_lst += lst

    all_ctr=Counter(big_lst)
    for s in stop_words: del all_ctr[s]
    logger.info("Common words which are possible candidates for new keywords:")
    logger.info(all_ctr.most_common(20))
    return

def dataframe_for_month(df, year=2017, month=1):
    """Yields a subset the dataframe with only those rows in the given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    if (mm==12): mm2,yyyy2=1,yyyy+1
    else: mm2,yyyy2=mm+1,yyyy

    logger.debug("Selecting rows for %i-%i" % (yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy,mm,1)) & (df.date<date(yyyy2,mm2,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def dataframe_for_6months(df, year=2017, month=1):
    """Yields a subset the dataframe with those rows for last 6 months up to given month
       Returns a dataframe
    """
    mm=month
    yyyy=year
    # calc month after this one, to set upper limit for search
    if (mm==12): mm_end,yyyy_end=1,yyyy+1
    else: mm_end,yyyy_end=mm+1,yyyy
    # calc 6 months ago, to set lower limit for search
    if (mm<6): mm_start,yyyy_start=mm+7,yyyy-1
    else: mm_start,yyyy_start=mm-5,yyyy

    logger.debug("Selecting rows for %i-%i to %i-%i" % (yyyy_start,mm_start,yyyy,mm))
    month_df = df.loc[ (df.date>=date(yyyy_start,mm_start,1)) & (df.date<date(yyyy_end,mm_end,1)) ]
    logger.debug("Found %i rows for this month" % len(month_df.index))

    return month_df

def found_word_list(df_col_list, kwd):
    """Passed a set of dataframe columns (i.e, a set of Series), together with a keyword
       Return a list (actually a Series) with True wherever at least one of the dataframe columns
       contains kwd or a synonym in the relevant columns, and False otherwise.
       Can be used to subset the original dataframe to pick out the rows with the keyword in them:
           df[found_word_list({df.wtlma, df.actions},"foo")]   => subset of df with "foo" in one of the columns
    """
    found_list = [0] * df_col_list[0].count()   # set up with list of zeros, same length as first column in the list
    for col in df_col_list:
        found_list = bool_list_of_occurrences(col, kwd) | found_list

    return found_list

def kwds_in_wtlma_actions(df,keyword_list):
    """ Count each keyword in the dataframe's important columns, wtlma and actions.
        Returns a Counter collection of (keyword:count)
    """
    counter_words_to_freq = Counter()
    for w in keyword_list:
         counter_words_to_freq[w] += sum(found_word_list([df.wtlma, df.actions],w))

    return counter_words_to_freq
    
def kwds_in_objectives(df,keyword_list):
    """ Count each keyword in the dataframe's 'objectives' column.
        Returns a Counter collection of (keyword:count)
    """
    counter_words_to_freq = Counter()
    for w in keyword_list:
         counter_words_to_freq[w] += sum(found_word_list([df.objectives],w))

    return counter_words_to_freq


def count_rows_with_comments(df):
    """ Count the number of rows in dataframe with a comment in either <Want to Learn More About> or <Action Items>
        Returns the count
    """
    return (df["Want to Learn More About"].notnull() | df["Action Items"].notnull()).sum()

def counts_by_centre(dataframe):
    centre_counts = dataframe.Ctr.value_counts()
    counts_list = (centre_counts.PA if hasattr(centre_counts,"PA") else 0,
                   centre_counts.H if hasattr(centre_counts,"H") else 0,
                   centre_counts.NY1 if hasattr(centre_counts,"NY1") else 0,
                   centre_counts.LON1 if hasattr(centre_counts,"LON1") else 0,
                   centre_counts.SNG if hasattr(centre_counts,"SNG") else 0
                   )
    return counts_list

class SimpleGroupedColorFunc(object):
    """Create a color function object which assigns EXACT colors
       to certain words based on the color to words mapping

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        self.word_to_color = {word: color
                              for (color, words) in color_to_words.items()
                              for word in words}

        self.default_color = default_color

    def __call__(self, word, **kwargs):
        return self.word_to_color.get(word, self.default_color)

class GroupedColorFunc(object):
    """Create a color function object which assigns DIFFERENT SHADES of
       specified colors to certain words based on the color to words mapping.

       Uses wordcloud.get_single_color_func

       Parameters
       ----------
       color_to_words : dict(str -> list(str))
         A dictionary that maps a color to the list of words.

       default_color : str
         Color that will be assigned to a word that's not a member
         of any value from color_to_words.
    """

    def __init__(self, color_to_words, default_color):
        from wordcloud import get_single_color_func
        self.color_func_to_words = [
            (get_single_color_func(color), set(words))
            for (color, words) in color_to_words.items()]

        self.default_color_func = get_single_color_func(default_color)

    def get_color_func(self, word):
        """Returns a single_color_func associated with the word"""
        try:
            color_func = next(
                color_func for (color_func, words) in self.color_func_to_words
                if word in words)
        except StopIteration:
            color_func = self.default_color_func

        return color_func

    def __call__(self, word, **kwargs):
        return self.get_color_func(word)(word, **kwargs)

def file_wordcloud_for_month(keywords_for_month, useful_rows_for_month, year=2017, month=8):
    """Input is a keywords_for_month, a Counter, plus useful_rows_for_month, an integer, and year and month as integers.  
       Returns the filename where the wordcloud is stored
    """
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    # (to install, run "pip install wordcloud")
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt
    import os

    # set font so 30% occurrence of top word uses 196 point font
    percent = round(100*(keywords_for_month.most_common(1)[0][1]/useful_rows_for_month))
    font_for_biggest_word = round( 196 * percent/30 )
    logger.debug("font for word {} is {} based on {}% = {} / {}".format(
                               keywords_for_month.most_common(1)[0][0],
                                     font_for_biggest_word,
                                             percent,
                                                keywords_for_month.most_common(1)[0][1],
                                                        useful_rows_for_month))

    font_path = os.path.join("C:",os.sep,"Windows",os.sep,"Fonts",os.sep,'arial.ttf')
    logger.debug("Generating the wordcloud")
    keywords_for_month += Counter()    # remove any zero or negative counts from the list
    wc_for_month = WordCloud(font_path=font_path,
                             width=2500,height=500,
                             prefer_horizontal=1.0,
                             relative_scaling=0.7,
                             max_font_size=font_for_biggest_word,
                             background_color="white",
                             random_state=1
                             ).generate_from_frequencies(keywords_for_month)


    # Words that are not in any of the dict_colour_of_keywords values
    # will be colored with a grey single color function
    default_color = 'grey'

    # Apply our color function
    #grouped_color_func = GroupedColorFunc(dict_colour_of_keywords, default_color)
    grouped_color_func = SimpleGroupedColorFunc(dict_colour_of_keywords, default_color)
    wc_for_month.recolor(color_func=grouped_color_func)

    # Build the generated image
    plt.imshow(wc_for_month,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    #plt.show()

    filename = tmpdir+"wordcloud-"+calendar.month_name[month]+".png"
    plt.imsave(filename,wc_for_month,format="png")
    plt.close()

    return filename


def file_graph_for_month_kwd(kwd,kwd_pos,vals,months,line_color):
    fig,ax=plt.subplots(figsize=(5.75,3.25))
    ##############ax.clear()

    ax.plot(vals,line_color,linewidth=3)
    ax.set_axis_off()
    plt.ylim(min(vals)-0.05,max(vals)+0.05)
    m_minus_2_percent = "{0:.0f}%".format(vals[0] * 100)
    m_this_percent    = "{0:.0f}%".format(vals[2] * 100)
    ax.text(0,vals[0]+0.02,calendar.month_abbr[months[0]]+"\n"+m_minus_2_percent,fontsize=30)
    ax.text(2,vals[2]+0.02,calendar.month_abbr[months[2]]+"\n"+m_this_percent,fontsize=30)
    filename = tmpdir+"graph-"+str(kwd_pos)+".png"
    logger.debug("Saving %s graph for keyword %s in file %s" % (kwd_pos,kwd,filename))
    fig.savefig(filename,bbox_inches="tight")
    plt.close(fig)

    return filename

def file_donut_pie_for_month(values,name):
    fig, ax = plt.subplots()
    ax.axis('equal')
    outside, _ = ax.pie(values,startangle=90,counterclock=False,
                        colors=list(colour_list))
    ax.legend(("Palo Alto","Houston","NY","London","Singapore"),fontsize=24,bbox_to_anchor=(0.8,1.0),frameon=False)
    width = 0.50  #determines the thickness of donut rim
    plt.setp( outside, width=width, edgecolor='white')

    filename = tmpdir+"donut-"+re.sub(r"[& ]","_",str(name))+".png"
    logger.debug("Saving <%s> donut in file %s with values %r" % (name, filename, values))
    fig.savefig(filename,bbox_inches="tight")

    plt.close(fig)

    return filename

def donut_pie_for_industries(industries,slide_shapes):
    """Given a Series resulting from value_counts(), and a slide_shapes, we create
       the doughnut pie chart corresponding to the value counts on the slide_shapes object.
    """
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from collections import Counter

    if (JapanAndChinaToOther and 'Japan' in industries):
        industries.Other += industries.Japan
        del industries['Japan']

    if (JapanAndChinaToOther and 'China' in industries):
        industries.Other += industries.China
        del industries['China']

    industry_counts = Counter( dict(zip(industries.keys().tolist(),industries.tolist())) )    # put in Counter so we can retrieve in descending order

    logger.debug("Adding pie for industries with data: "+str(industry_counts))
    chart_data=ChartData()
    # This next is a bit obscure, but if i_c=Counter([('A':5),('B':4),('C':3)])
    # then zip(*i_c.most_common()) gives us a zip object which iterates out to a list of the keys, and the values, in descending order
    ind_lists = [i for i in zip(*industry_counts.most_common())] 
    chart_data.categories=ind_lists[0]
    chart_data.add_series("Industries",ind_lists[1])
    left_x=Mm(19)
    top_y=Mm(53)
    width=Mm(83); height=Mm(110)
    chart = slide_shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left_x,top_y, width, height, chart_data).chart   
    chart.has_title = False   #TODO: work out why we are still getting a title for the chart, even when we ask not to have one!!
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM 
    chart.legend.include_in_layout = 0 
    chart.legend.font.size = Pt(11) 

    return 

def donut_pie_for_centres(value_tuple,legend,slide_shapes,left_x,top_y,width,height):
    """Given a tuple of values, and a list of legend keys, such as (3,4,5,1,2) and ['EBC,'SNG','NY','H','LON']
       put a donut pie chart on the slide_shapes using the values from the list, in the given position & size, 
       with a legend underneath the pie of the data element names (the keys in the list).
       We use a list, not a dictionary, because we want to preserve order of the keys.
    """
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    logger.debug("Adding pie with data: "+str(value_tuple))
    chart_data=ChartData()
    chart_data.categories=legend
    chart_data.add_series("Centres",value_tuple)
    chart = slide_shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left_x,top_y, width, height, chart_data).chart   
    chart.has_title = False   #TODO: work out why we are still getting a title for the chart, even when we ask not to have one!!
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM 
    chart.legend.include_in_layout = 0 
    chart.legend.font.size = Pt(10) 

    return 

def write_customer_list(df_for_kwd,text_frame):
    text_frame.text = "Briefings"
    font=text_frame.paragraphs[0].font
    font.size = Pt(12)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for i in (df_for_kwd)["Account Name"].str.title().iteritems():
        new_para = text_frame.add_paragraph()
        font=new_para.font
        font.size = Pt(8)
        font.color.theme_color = MSO_THEME_COLOR.DARK_2
        new_para.text=i[1]

    return

def write_top_keywords(text_frame,titleText,kwd_counts_for_ind,rowcount,percent=True,cutoff=0):
    """Write a section headed by titleText, with bullets of the top 4 elements of Counter
       kwd_counts_for_ind plus their percentage of use (if percent=True), where
       the percentage is calculated with rowcount as the denominator.  Stop before top
       4 elements if they start to fall below the cutoff percentage.
         e.g:
           c=Counter([("alpha":5),("bravo":4),("charlie":3),("delta":1))])
           write_top_keywords(tf,"My list",c,10,percent=True,cutoff=20)
         yields:
           My list
            - alpha - 50%
            - bravo - 40%
            - charlie - 30%
    """
    text_frame.text = titleText
    font=text_frame.paragraphs[0].font
    font.size = Pt(10)
    font.bold = True
    font.color.theme_color = MSO_THEME_COLOR.DARK_1

    for k,v in kwd_counts_for_ind.most_common(4):
        if (100*v/rowcount>=cutoff):
            new_para = text_frame.add_paragraph()
            font=new_para.font
            font.size = Pt(8)
            font.color.theme_color = MSO_THEME_COLOR.DARK_2
            if (percent==True):
                new_para.text=" - "+k+" - "+"{0:.0f}%".format(v/rowcount * 100)
            else:
                new_para.text=" - "+k

    return

def find_text_in_shapes(slide_shapes,searchword):
    """Find the given text in the list of powerpoint shapes passed in slide_shapes
    """
    found_idx = -1
    for i in range(len(slide_shapes)):
        if (slide_shapes[i].has_text_frame):
            text_frame = slide_shapes[i].text_frame
            if (text_frame.text.find(searchword)>=0):
                found_idx = i
                break

    return found_idx

def replace_text_in_shape(slide_shapes,find,use,slidename):
    """Find the placeholder text "find" in the slide_shapes, replace it with "use",
       reporting an error on "slidename" if we can't find "find"
       Returns the text_frame written
    """
    i = find_text_in_shapes(slide_shapes,find)
    if (i>=0):
        text_frame = slide_shapes[i].text_frame
        text_frame.text = use
    else:
        logger.error("Could not find %s placeholder on %s" % (find,slidename))

    return

def new_run_in_slide(para,text='text',fontname='Arial',fontsize=24):
    "Add new run to a paragraph in a powerpoint slide"
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.name = fontname
    return run

import struct
import imghdr

def get_image_size(fname):
    '''Determine the image type of fname and return its width, height.
    from draco (via stackoverflow)'''
    with open(fname, 'rb') as fhandle:
        head = fhandle.read(24)
        if len(head) != 24:
            return
        if imghdr.what(fname) == 'png':
            check = struct.unpack('>i', head[4:8])[0]
            if check != 0x0d0a1a0a:
                return
            width, height = struct.unpack('>ii', head[16:24])
        elif imghdr.what(fname) == 'gif':
            width, height = struct.unpack('<HH', head[6:10])
        elif imghdr.what(fname) == 'jpeg':
            try:
                fhandle.seek(0) # Read 0xff next
                size = 2
                ftype = 0
                while not 0xc0 <= ftype <= 0xcf:
                    fhandle.seek(size, 1)
                    byte = fhandle.read(1)
                    while ord(byte) == 0xff:
                        byte = fhandle.read(1)
                    ftype = ord(byte)
                    size = struct.unpack('>H', fhandle.read(2))[0] - 2
                # We are at a SOFn block
                fhandle.seek(1, 1)  # Skip `precision' byte.
                height, width = struct.unpack('>HH', fhandle.read(4))
            except Exception: #IGNORE:W0703
                return
        else:
            return
        return width, height

def add_icon(ss,icon_name,left,top,small=False):
    """Add the icon in icondir/icon_name to the slide shape ss at coordinates left,top
       If small is True, make it a smaller icon.
    """
    ## ok if icon exists, and use error icon if not
    logger.debug("adding icon <{}> at ({},{})".format(icon_name,left,top))
    try:
        filename = icondir+icon_name+".png"
        icon_width, _icon_height = get_image_size(filename)
        pixels_per_mm = 6
        if small:
            icon_width = 0.4*icon_width
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(12),Mm(12))
        else:
            ss.add_picture(filename,round(left-icon_width/pixels_per_mm),top,Mm(34),Mm(34))
    except:
        logger.warning("Could not find icon {}".format(icondir+icon_name+'.png'))
    return


def hpe_color_fn(word, font_size, position, orientation, random_state=None, **kwargs):
    """Returns colour to use for given word, font_size, etc.
    """
    import random
    colour = random.choice(colour_list)
    return colour

def file_wordcloud_for_partners(partner_counts):
    """Input is partner_counts, a Counter of partner names and occurrences
       Returns the filename where the wordcloud is stored
    """
    ##Build a wordcloud, using the wordcloud code from Andreas Mueller
    from wordcloud import WordCloud
    import matplotlib.pyplot as plt
    import os

    # Chop list after most common 30 partners (note, need to use dict() around most_common() as it returns a list)
    p_counts = Counter(dict(partner_counts.most_common(30)))
    # set font so 15% occurrence of top word uses 196 point font
    percent = round(100*(p_counts.most_common(1)[0][1]/sum(p_counts.values())))
    font_for_biggest_word = round( 196 * percent/15 )

    font_path = "C:"+os.sep+"Windows"+os.sep+"Fonts"+os.sep+'arial.ttf'
    logger.debug("Generating the wordcloud for Partners")
    p_counts += Counter()    # remove any zero or negative counts from the list
    wc_partners = WordCloud(font_path=font_path,
                            width=2000,height=500,
                            prefer_horizontal=1.0,
                            relative_scaling=0.7,
                            max_font_size=font_for_biggest_word,
                            background_color="white",
                            random_state=1
                            ).generate_from_frequencies(p_counts)

    # Re-colour to the required colour se
    wc_partners.recolor(color_func=hpe_color_fn, random_state=3)

    # Build the generated image
    plt.imshow(wc_partners,interpolation='bilinear')
    plt.axis("off")
    plt.figure()
    # Save the wordcloud to file
    filename = tmpdir+"partner_wordcloud.png"
    plt.imsave(filename,wc_partners,format="png")
    plt.close()
    return filename

def sentiment_by_month(df, count, year, month):
    # Calculate the sentiment by month for the last 'count' months from 'month' in 'year'
    # Returns an array with the average compound sentiment score per month, from
    # most recent month to oldest month (so July, June, May, in that order, etc)
    from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
    assert (count<12),"only works for intervals of up to 12 months"   
    analyzer = SentimentIntensityAnalyzer()
    result=[]
    for i in range(0,count):
        m = month-i
        y = year
        if (m<=0):   
            m += 12
            y -= 1
        month_df = dataframe_for_month(df, year=y, month=m)
        avg_snt = tot_snt = {'neg': 0.0, 'neu':0.0, 'pos': 0.0, 'compound':0.0}
        n=0
        for _, row in month_df.iterrows():
            cell = row["Customer Overall Comments"]
            if (isinstance(cell, str)):
                snt = analyzer.polarity_scores(cell)
                n += 1
                tot_snt['neg'] += snt['neg']
                tot_snt['neu'] += snt['neu']
                tot_snt['pos'] += snt['pos']
        if (n > 0):        
            avg_snt['neg'] = round(tot_snt['neg'] / n, 3)         
            avg_snt['neu'] = round(tot_snt['neu'] / n, 3)         
            avg_snt['pos'] = round(tot_snt['pos'] / n, 3)         
            avg_snt['compound'] = round(tot_snt['compound'] / n, 3)         
        logger.info("Avg sentiment (n={}), month {}: {}".format(n,m,str(avg_snt)))
        result.append(avg_snt['compound'])
    return result


def top_sentiment_in_month(month_df, count=4):
    # Return dictionary with top [count] comments for given month, by sentiment
    # Each dictionary item is   {"user comment in full": compound_sentiment_score}
    from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
    assert (isinstance(count,int)),"only works for values of count which are integers"
    assert (count>0),"only works for values of count which are greater than zero"
    analyzer = SentimentIntensityAnalyzer()
    results={}
    lowest_top_snt = 0.0  # lowest sentiment to make it into top N
    found_count=0
    for _, row in month_df.iterrows():
        cell = row["Customer Overall Comments"]
        if (isinstance(cell, str)):
            this_snt = analyzer.polarity_scores(cell)
            if (found_count < count):
                # Haven't yet found [count] items to return, so this will automatically be added
                found_count += 1 
                results[cell] = this_snt['compound']
            elif (this_snt["compound"] > lowest_top_snt):
                # This item is more positive than lowest in top so far
                # So find lowest scoring item already saved, and discard 
                del results[ min(results, key=results.get) ]
                results[cell]=this_snt['compound']
            lowest_top_snt = min(results.values())
    return results


logger.info("Importing excel file "+excel_file+" with header in row "+str(header_row))
# Note that "header" is zero-indexed, so must subtract one from header_row
all_df = pd.read_excel(open(excel_file,'rb'),header=header_row-1,usecols="A:S")   

logger.debug("Adding structured columns")
all_df.insert(loc=0,column='date',value=all_df['Visit Date'].apply(make_date))
all_df.insert(loc=1,column='wtlma',value=all_df['Want to Learn More About'].apply(tidy_text))
all_df.insert(loc=2,column='actions',value=all_df['Action Items'].apply(tidy_text))
all_df.insert(loc=3,column='objectives',value=all_df['Objectives'].apply(tidy_text))

print_new_candidate_words(all_df,stop_words,top_n=40)
if stop_after_wordcheck:
    sys.exit()


## Given the month and year, calc the number of the previous few months
if (mm==1): mm_minus_1,year_for_mm_minus_1 = 12, yyyy-1
else: mm_minus_1, year_for_mm_minus_1 = mm-1, yyyy
if (mm<=2): mm_minus_2,year_for_mm_minus_2 = mm+10, yyyy-1
else: mm_minus_2, year_for_mm_minus_2 = mm-2, yyyy
if (mm<6): mm_6_before,year_for_mm_6_before = mm+7, yyyy-1   #6 months before Jan is Aug
else: mm_6_before, year_for_mm_6_before = mm-5, yyyy    #6 months before Jul is Feb

### Now start to generate the powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt, Mm

# These appear to be the layouts used in the master for this slide deck
LAYOUT_TITLE_WITH_DARK_PICTURE = 0
LAYOUT_TITLE_SLIDE_WITH_NAME   = 1
LAYOUT_DIVIDER                 = 2
LAYOUT_TITLE_ONLY              = 3
LAYOUT_TITLE_AND_SUBTITLE      = 4
LAYOUT_BLANK                   = 5

## Open up the source presentation
prs = Presentation('GCA_Customer_Insights_Month-Year.pptx')
this_month = calendar.month_name[mm]
earliest_month = calendar.month_name[mm_minus_2]
logger.info("Creating presentation for %s" % (this_month))

################################################
## Modify existing title slide with month
################################################
logger.info("Modifying text in slide 0")
s = prs.slides[0]
slide_shapes = s.shapes
text_frame = slide_shapes[0].text_frame  # should be the the title textframe
#clear existing text, and write new text into title textframe
text_frame.clear()
# First para with a run of 60-point Arial font
new_run_in_slide(text_frame.paragraphs[0],text='Customer Insights',fontname='Arial',fontsize=60)
# Second para with a run of 32-point font
new_run_in_slide(text_frame.add_paragraph(),text='Learnings from '+this_month+' EBC/CEC visits',fontsize=32)

################################################
## 3rd slide: count the keywords for this month and build the wordcloud
################################################
logger.info("3rd slide: large wordcloud for this month, and top 3 keywords")
df_for_month = dataframe_for_month(all_df, year=yyyy, month=mm)
kwd_count_for_month = kwds_in_wtlma_actions(df_for_month,vocab)
useful_rows_in_m = count_rows_with_comments(df_for_month)
logger.info("Top keyword/counts for month %i : %r" % (mm,kwd_count_for_month.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_month, useful_rows_in_m,
                         year=yyyy,month=mm)
## Modifying main wordcloud slide by changing title and adding pic for this month's wordcloud
logger.debug("Modifying text and adding wordcloud in slide 2")
s = prs.slides[2]
slide_shapes=s.shapes
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="In "+this_month+" customers wanted to learn more about...",
       fontname="Arial",fontsize=28)
left=Mm(12.5); top=Mm(100)
slide_shapes.add_picture(tmpdir+"wordcloud-"+this_month+".png",left,top,height=Mm(63))

##Find where the placeholders are for the top 3 keywords, update them, then add their icons
top_3 = kwd_count_for_month.most_common(3)   # top 3 keywords for most recent month in list with their counts
kwd0=top_3[0][0]
kwd1=top_3[1][0]
kwd2=top_3[2][0]
logger.debug("Top 3 keywords for current month are %s %s %s" % (kwd0,kwd1,kwd2))
replace_text_in_shape(slide_shapes,find="topword_1",use=kwd0,slidename="3rd slide")
replace_text_in_shape(slide_shapes,find="topword_2",use=kwd1,slidename="3rd slide")
replace_text_in_shape(slide_shapes,find="topword_3",use=kwd2,slidename="3rd slide")
t=Mm(51)
add_icon(slide_shapes,kwd0,left=Mm(47),top=t)
add_icon(slide_shapes,kwd1,left=Mm(152),top=t)
add_icon(slide_shapes,kwd2,left=Mm(255),top=t)

## @15mar18: add text into notes for this slide to show actual counts for top 10 words
notes_for_slide = s.notes_slide
notes_tf = notes_for_slide.notes_text_frame
notes_tf.text = ("Found %i rows with comments in this month.\nTop ten keyword/counts for %s: \n%r" % (useful_rows_in_m,this_month,kwd_count_for_month.most_common(10)) )

################################################
## 4th slide: count the keywords for previous two months, and build their wordclouds.
################################################
logger.info(">>>> 4th slide: three wordclouds for most recent 3 months")
df_for_month_minus_1 = dataframe_for_month(all_df, year=year_for_mm_minus_1, month=mm_minus_1)
kwd_count_for_m_minus_1 = kwds_in_wtlma_actions(df_for_month_minus_1,vocab)
logger.info("Top keyword/counts for month %i : %r" % (mm_minus_1,kwd_count_for_m_minus_1.most_common(5)) )

df_for_month_minus_2 = dataframe_for_month(all_df, year=year_for_mm_minus_2, month=mm_minus_2)
kwd_count_for_m_minus_2 = kwds_in_wtlma_actions(df_for_month_minus_2,vocab)
logger.info("Top keyword/counts for month %i : %r" % (mm_minus_2,kwd_count_for_m_minus_2.most_common(5)) )

useful_rows_in_m_2 = count_rows_with_comments(df_for_month_minus_2)
useful_rows_in_m_1 = count_rows_with_comments(df_for_month_minus_1)
logger.info("Number of useful rows in months -2,-1,0 are %i, %i, %i" % (useful_rows_in_m_2,useful_rows_in_m_1,useful_rows_in_m))

file_wordcloud_for_month(kwd_count_for_m_minus_1, useful_rows_in_m_1, year=year_for_mm_minus_1,month=mm_minus_1)
file_wordcloud_for_month(kwd_count_for_m_minus_2, useful_rows_in_m_2, year=year_for_mm_minus_2,month=mm_minus_2)

## Modifying 3-month wordcloud slide by adding pic for last 3 months' wordcloud
logger.debug("Adding three wordclouds in 4th slide")
s = prs.slides[3]
slide_shapes=s.shapes
left=Mm(93)
top=Mm(36)
h=Mm(38)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_2]+".png",
                         left,top,height=h)
top=Mm(83)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_1]+".png",
                         left,top,height=h)
top=Mm(128)
slide_shapes.add_picture(tmpdir+"wordcloud-"+this_month+".png",
                         left,top,height=h)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
replace_text_in_shape(slide_shapes,find="Month-2",use=calendar.month_name[mm_minus_2],slidename="4th slide")
replace_text_in_shape(slide_shapes,find="Month-1",use=calendar.month_name[mm_minus_1],slidename="4th slide")
replace_text_in_shape(slide_shapes,find="Month-0",use=calendar.month_name[mm],slidename="4th slide")

################################################
## 5th slide: for top 3 keywords for this month graph their usage
################################################
logger.info(">>>> 5th slide: top keywords for this month")
months=[mm_minus_2,mm_minus_1,mm]

kwd0_c2 = kwd_count_for_m_minus_2[kwd0] if (kwd0 in kwd_count_for_m_minus_2) else 0
kwd0_c1 = kwd_count_for_m_minus_1[kwd0] if (kwd0 in kwd_count_for_m_minus_1) else 0
kwd0_c0 = kwd_count_for_month[kwd0]     # must have keyword as it came from this dictionary
vals_kwd0=[kwd0_c2/useful_rows_in_m_2, kwd0_c1/useful_rows_in_m_1, kwd0_c0/useful_rows_in_m]
file_linegraph_topic1 = file_graph_for_month_kwd(kwd0,"1st",vals_kwd0,months,colour_list[0])
logger.debug("Kwd0 is %s, data %r" % (kwd0,vals_kwd0))
kwd1_c2 = kwd_count_for_m_minus_2[kwd1] if (kwd1 in kwd_count_for_m_minus_2) else 0

kwd1_c1 = kwd_count_for_m_minus_1[kwd1] if (kwd1 in kwd_count_for_m_minus_1) else 0
kwd1_c0 = kwd_count_for_month[kwd1]     # must have keyword as it came from this dictionary
vals_kwd1=[kwd1_c2/useful_rows_in_m_2, kwd1_c1/useful_rows_in_m_1, kwd1_c0/useful_rows_in_m]
file_linegraph_topic2 = file_graph_for_month_kwd(kwd1,"2nd",vals_kwd1,months,colour_list[1])
logger.debug("Kwd1 is %s, data %r" % (kwd1,vals_kwd1))

kwd2_c2 = kwd_count_for_m_minus_2[kwd2] if (kwd2 in kwd_count_for_m_minus_2) else 0
kwd2_c1 = kwd_count_for_m_minus_1[kwd2] if (kwd2 in kwd_count_for_m_minus_1) else 0
kwd2_c0 = kwd_count_for_month[kwd2]     # must have keyword as it came from this dictionary
vals_kwd2=[kwd2_c2/useful_rows_in_m_2, kwd2_c1/useful_rows_in_m_1, kwd2_c0/useful_rows_in_m]
file_linegraph_topic3 = file_graph_for_month_kwd(kwd2,"3rd",vals_kwd2,months,colour_list[2])
logger.debug("Kwd0 is %s, data %r" % (kwd2,vals_kwd2))

## Build a subset of the dataframe for last 3 months that uses each of the top 3 kwds in this month
df_for_3months = pd.concat([df_for_month,df_for_month_minus_1,df_for_month_minus_2])
col_list = [df_for_3months.actions, df_for_3months.wtlma]  # list of columns to search for kwds
df_for_kwd0 = df_for_3months.loc[found_word_list(col_list,kwd0)]
df_for_kwd1 = df_for_3months.loc[found_word_list(col_list,kwd1)]
df_for_kwd2 = df_for_3months.loc[found_word_list(col_list,kwd2)]

## Create donut pies showing split of visits expressing interest in top 3 topics by centre over last 3 months
kwd0_counts = counts_by_centre(df_for_kwd0)
file_donut_topic1 = file_donut_pie_for_month(counts_by_centre(df_for_kwd0),"1st")
file_donut_topic2 = file_donut_pie_for_month(counts_by_centre(df_for_kwd1),"2nd")
file_donut_topic3 = file_donut_pie_for_month(counts_by_centre(df_for_kwd2),"3rd")

## Add the line graphs and donut pies to the Top 3 Customer Interests chart
logger.debug("Adding line graphs, donuts and customers to the Top 3 Customer Interests slide (5th slide)")
s = prs.slides[4]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Top 3 Customer Interests: "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
replace_text_in_shape(slide_shapes,find="Topic1",use=kwd0,slidename="5th slide")
replace_text_in_shape(slide_shapes,find="Topic2",use=kwd1,slidename="5th slide")
replace_text_in_shape(slide_shapes,find="Topic3",use=kwd2,slidename="5th slide")

#Add the line graphs and the donuts for each of the topics
top=Mm(46); h=Mm(25); w=Mm(43)
slide_shapes.add_picture(file_linegraph_topic1,Mm(19),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic2,Mm(120),top,height=h,width=w)
slide_shapes.add_picture(file_linegraph_topic3,Mm(223),top,height=h,width=w)
top=Mm(43); h=Mm(28); w=Mm(57)
slide_shapes.add_picture(file_donut_topic1,Mm(58),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic2,Mm(160),top,height=h,width=w)
slide_shapes.add_picture(file_donut_topic3,Mm(263),top,height=h,width=w)

#Find where the placeholders are for the customer lists and update them
col_list = [df_for_month.actions, df_for_month.wtlma]  # list of columns to search for kwds
df_for_kwd0_month0 = df_for_month.loc[found_word_list(col_list,kwd0)]
df_for_kwd1_month0 = df_for_month.loc[found_word_list(col_list,kwd1)]
df_for_kwd2_month0 = df_for_month.loc[found_word_list(col_list,kwd2)]
i = find_text_in_shapes(slide_shapes,"Customers1")
if (i>=0):
    logger.debug("Writing list of %i customers for first keyword" % (len(df_for_kwd0)))
    write_customer_list(df_for_kwd0_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers1 placeholder")

i = find_text_in_shapes(slide_shapes,"Customers2")
if (i>=0):
    logger.debug("Writing list of %i customers for second keyword" % (len(df_for_kwd1)))
    write_customer_list(df_for_kwd1_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers2 placeholder")

i = find_text_in_shapes(slide_shapes,"Customers3")
if (i>=0):
    logger.debug("Writing list of %i customers for third keyword" % (len(df_for_kwd2)))
    write_customer_list(df_for_kwd2_month0,slide_shapes[i].text_frame)
else:
    logger.error("Could not find Customers3 placeholder")

################################################
## 9th slide: Now generate the Industry Insights donuts and top keyword lists
################################################
logger.info(">>>> 9th slide: industries and how they show up across the centres")

## Sort through the dataframes, extracting relevant information, saving it in dictionaries for later use
df_for_ind={}
kwd_counts_for_ind={}
ctr_counts_for_ind={}

logger.debug("Generating dataframes for each industry for last 3 months")
for ind in industry_list:
    df_for_ind[ind] = df_for_3months[df_for_3months["Industry"]==ind]
    ctr_counts_for_ind[ind]=counts_by_centre( df_for_ind[ind] )
    kwd_counts_for_ind[ind] = kwds_in_wtlma_actions(df_for_ind[ind],vocab)
    logger.debug("Top keyword/counts for %s: %r" % (ind,kwd_counts_for_ind[ind].most_common(4)) )

## Now write out the charts and text on industry slide
s = prs.slides[8]
slide_shapes=s.shapes

#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Industry Insights "+earliest_month+"-"+this_month,
       fontname="Arial",fontsize=28)

#Add the one big donut showing volumes for each industry to the slide
logger.debug("Adding main donut to Industry Insights slide (9th slide)")
industry_counts = df_for_3months["Industry"].value_counts()
donut_pie_for_industries(industry_counts,slide_shapes)

#Add the individual industry donuts broken down by centre to the slide
#Need to pick the top N, excluding "Other", from the list in industry_list
logger.debug("Adding donuts for top industries in each centre (9th slide)")
pie_left=(Mm(110),Mm(164),Mm(217),Mm(272), Mm(110),Mm(164),Mm(217),Mm(272))
pie_top =(Mm(40), Mm(40), Mm(40), Mm(40),  Mm(115),Mm(115),Mm(115),Mm(115))
pie_h=Mm(50); pie_w=Mm(35)
icon_left=(Mm(143),Mm(198),Mm(252),Mm(306))    # Placement for icons in each of 4 columns
icon_top=(Mm(48),Mm(65),Mm(82),Mm(122),Mm(138),Mm(156))   # Placement of icons in two sets of rows, three in each row
for n,ind in enumerate(['Public Sector','Fin Svcs','RCG','Energy','Health & LS','Mfg','CME','Travel & Trans']):
    #['Public Sector','Fin Svcs','RCG','Energy','Health & LS','Mfg','CME','Travel & Trans']
    logger.debug("Writing %s as industry %i" %(ind,n))
    #Write the industry name as the main title for this box
    replace_text_in_shape(slide_shapes,"Industry-{}".format(n),ind,"9th slide")
    #Add the donut showing breakdown of centres that hosted this industry
    donut_pie_for_centres(ctr_counts_for_ind[ind],['PA','H','NY','L','SNG'],slide_shapes,pie_left[n],pie_top[n],pie_w,pie_h)
    #Now write the list of top interests for this industry
    for interest_idx, p in enumerate(kwd_counts_for_ind[ind].most_common(3)):
        interest=p[0]
        logger.debug("Replacing {}-interest-{} with {} and its icon".format(n,interest_idx,interest))
        replace_text_in_shape(slide_shapes,"{}-interest-{}".format(n,interest_idx),interest,"9th slide")
        add_icon(slide_shapes,interest,left=icon_left[n%4],top=icon_top[3*(n//4)+interest_idx],small=True)

################################################
## 10th slide: Partner insights
################################################
logger.info(">>>> 10th slide: top partner keywords and partner attendance broken out by centre")
# rather copmlex expression to find which rows are partners who are attending as partners, or partner-led briefings for customers
df_for_partners = df_for_3months.loc[ ( df_for_3months['Account Type'].isin(["Channel/ Reseller","Systems Integrator"])
                                        & (df_for_3months['Partner / Customer']!="Customer")
                                      ) | (df_for_3months['Partner / Customer']=="Partner")
                                    ]
kwd_count_for_partners = kwds_in_wtlma_actions(df_for_partners,vocab)
useful_rows_in_partners = count_rows_with_comments(df_for_partners)

df_channel = df_for_partners[ df_for_partners['Account Type']=="Channel/ Reseller" ]
df_SI = df_for_partners[ df_for_partners['Account Type']=="Systems Integrator" ]
df_accompanied = df_for_partners[ df_for_partners['Partner / Customer']=="Partner" ]


SI_count=[len(df_SI[df_SI.Ctr==c]) for c in centres]
Channel_count=[len(df_channel[df_channel.Ctr==c]) for c in centres]
Attended_count=[len(df_accompanied[df_accompanied.Ctr==c]) for c in centres]
logger.debug("Centres being looked at: %r" %(centres))
logger.debug("Volume by centre - partner: %r" %(Channel_count))
logger.debug("Volume by centre - SI: %r" %(SI_count))
logger.debug("Volume by centre - attended: %r" %(Attended_count))

# @edit requested by Tina 26th Feb: combine SI with Channel/Reseller
channel_plus_SI=[x+y for x,y in zip(Channel_count, SI_count)]

## Start to update the slide
s = prs.slides[9]
slide_shapes=s.shapes

#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Partner Insights ("+earliest_month+"-"+this_month+")",
       fontname="Arial",fontsize=28)

# Update the top 3 most common keywords, and their percentages
logger.debug("top 3 partner interests: %r" % (kwd_count_for_partners.most_common(4)))
t=Mm(66)
l=[Mm(32),Mm(69),Mm(106)]
for n,p in enumerate(kwd_count_for_partners.most_common(3)):
    # p is (keyword: count) for each of the top most common keywords
    replace_text_in_shape(slide_shapes,"Interest-{}".format(n),p[0],"10th slide")
    add_icon(slide_shapes,p[0],top=t,left=l[n],small=True)
    replace_text_in_shape(slide_shapes,"Score-{}".format(n),"{0:.0f}%".format(100*p[1]/useful_rows_in_partners),"10th slide")

# @add: show pie charts and partner wordcloud - 19sep18
# @del: remove pie charts! - 21sep18
partner_file = "Visits with Partners.xlsx"
logger.info("Importing excel file for partners: "+partner_file+", with header in row "+str(header_row))
# Read the dataframe of all the names, partner names, dates, etc, of the individual people from partners who attended a centre
partner_df = pd.read_excel(open(partner_file,'rb'),header=0,usecols="A:H")   
# Subset the dataframe to the 3 month window we are intested in
logger.debug("Selecting rows from Partner file from %i-%i to %i-%i" % (year_for_mm_minus_2,mm_minus_2,yyyy,mm))
start_date = pd.Timestamp(year_for_mm_minus_2,mm_minus_2,1)
if (mm==12): # if current month is december, end date is 1st Jan nest year
    end_date = pd.Timestamp(yyyy+1,1,1)
else: 
    end_date = pd.Timestamp(yyyy,mm+1,1)
partner_3m_df = partner_df.loc[ (partner_df['Visit: Arrival Date']>= start_date) & (partner_df['Visit: Arrival Date'] < end_date) ]
## Now generate wordcloud, based on the dataframe we just created
grp_by_partner = partner_3m_df.groupby(['Attendee Company Name']) # group has one row per attendee, grouped by partner name
ptr_counts = grp_by_partner.size() # counts number of each partner - need to turn this into a python Counter() object
ptr_counter=Counter()
for i,v in ptr_counts.iteritems():
    ptr_counter[i]=v
ptr_counter['HPE']=0   # Zap out the entry (if any) for HPE as a parter - due to "garbage in" from briefing mgrs
file_wc_ptr=file_wordcloud_for_partners(ptr_counter)
# Add it to the slide template
slide_shapes.add_picture(file_wc_ptr, Mm(142),Mm(50), height=Mm(40),width=Mm(178))

################################################
## 11th slide: Top interests and industries for last 6 months
################################################
logger.info(">>>> 11th slide: top 5 interests, top 3 industries, and their top interests, by centre, for last 6 months")
df_6months = dataframe_for_6months(all_df, year=yyyy, month=mm)

# Calculate the sentiment by month for the last 6 months
sentiment_6months = sentiment_by_month(df_6months, count=6, year=yyyy, month=mm)
if sentimentCalcs:
    # Get top scoring comments for this month
    top_comments_in_month = top_sentiment_in_month(df_for_month,count=4)   
    for c in top_comments_in_month:
        print(top_comments_in_month[c],":",c)
    
#build a dictionary of dataframes subsetted by centre, a dictionary of keywords by centre,
#and a dict of top 3 industries per centre and their keywords (where the key is a tuple of (centre, industry) )
dfs_6m_ctr={}
kwd_counts_6m_ctr={}
industry_counts_6m = df_6months["Industry"].value_counts()
kwd_counts_6m_for_top_inds={}
commented_rows_for_6m={}

for c in centres:
    logger.info("Working on counts for {}".format(c))
    dfs_6m_ctr[c]=df_6months[df_6months.Ctr==c]
    #First, the top keywords for that Centre
    kwd_counts_6m_ctr[c]=kwds_in_wtlma_actions(dfs_6m_ctr[c],vocab)
    logger.debug("Top keyword/counts %s: %r" % (c,kwd_counts_6m_ctr[c].most_common(5)) )
    #Now the top keywords in each industry for that centre
    industry_counts_6m[c] = (dfs_6m_ctr[c])["Industry"].value_counts()
    for ind in (industry_counts_6m[c]).index:
        this_df = dfs_6m_ctr[c][ dfs_6m_ctr[c]["Industry"]==ind ]
        kwd_counts_6m_for_top_inds[(c,ind)]=kwds_in_wtlma_actions(this_df, vocab)
        commented_rows_for_6m[(c,ind)]=count_rows_with_comments(this_df)
        logger.debug("Top keyword/counts in %s for %s: %r" % (c,ind,kwd_counts_6m_for_top_inds[(c,ind)].most_common(3)) )


## Build the slide: add the interests, industries, and per-industry interests, for each of the centres
logger.debug("Setting the title ")
s = prs.slides[10]
slide_shapes=s.shapes
#Update the title
title_frame = slide_shapes.title.text_frame
title_frame.clear()
new_run_in_slide(title_frame.paragraphs[0],text="Breakdown by centre ("+calendar.month_name[mm_6_before]+"-"+this_month+")",
       fontname="Arial",fontsize=28)

#Update, by centre, the top interests, and the top industries with their top interests
logger.info("For each centre, getting top 5 interests, and top 3 industries and their interests.")
top_pos=[Mm(48),Mm(73),Mm(102),Mm(127),Mm(152)]   #distances to top of icon for each row
left_pos=[Mm(35),Mm(65),Mm(95),Mm(125),Mm(154)]   #distances to centre of icon for each row
for row,ctr in enumerate(["PA","H","NY1","LON1","SNG"]):  #iterate the centres in the order they appear on the slide
    #First, the top interests for this centre
    for col,p in enumerate(kwd_counts_6m_ctr[ctr].most_common(5)):
        # p is (keyword: count) for each of the keywords, so p[0] is the keyword itself
        replace_text_in_shape(slide_shapes,"{}-interest-{}".format(ctr,col),p[0],"11th slide")
        add_icon(slide_shapes,p[0],left=left_pos[col],top=top_pos[row],small=True)
    #Next, the top industries with their interests for this centre - industry_counts_6m[c] is already ordered highest->lowest count
    n=0  #count how many displayed - need to do this separately from the loop count, as we ignore "Other" as an industry group
    for ind in industry_counts_6m[ctr].index:
        if   (ind!="Other"):
            logger.debug("For centre <%s> industry <%i> is <%s>" %(ctr,n,ind))
            #Write the list of top interests for this industry
            idx = find_text_in_shapes(slide_shapes,"{}-industry-{}".format(ctr,n))
            if (idx>=0):
                write_top_keywords(slide_shapes[idx].text_frame,
                                   ind,
                                   kwd_counts_6m_for_top_inds[(ctr,ind)],
                                   commented_rows_for_6m[(ctr,ind)],
                                   percent=False,
                                   cutoff=20
                                  )
            else:
                logger.error("Could not find <{}-industry-{}> placeholder on 11th slide - idx={}".format(ctr,n,idx))
            n+=1   #increment the number of industries written out for this centre
        if (n>=3): break   #exit the loop after writing in 3 industries+interests

## @24apr18: add text into notes for this slide to show actual counts for top words per centre
notes_for_slide = s.notes_slide
notes_tf = notes_for_slide.notes_text_frame
notes_tf.text = "Counts by centre for last 6 months:\n"
for row,ctr in enumerate(["PA","H","NY1","LON1","SNG"]):  #iterate the centres in the order they appear on the slide
    notes_tf.text += "For %s, top eight items were: %r\n" % (ctr, kwd_counts_6m_ctr[ctr].most_common(8))

################################################
## 13th slide: EBC specific volumes & interests for last 6 months
################################################
if False:
    logger.info(">>>> 13th slide: EBC specific industry volumes and top interests")
    ## Generate the image for the big donut showing volume for all industries across the months
    PA_industry_counts = dfs_6m_ctr["PA"]["Industry"].value_counts()
    file_donut_PA_ind_vols = file_donut_pie_for_industries(PA_industry_counts,center="PA")

    ## Build the slide
    logger.debug("Setting the title ")
    s = prs.slides[12]
    slide_shapes=s.shapes
    #Update the title
    title_frame = slide_shapes.title.text_frame
    title_frame.clear()
    new_run_in_slide(title_frame.paragraphs[0],text="EBC six month view ("+calendar.month_name[mm_6_before]+"-"+this_month+")",
        fontname="Arial",fontsize=28)

    #Add the one big donut showing volumes for each industry to the slide
    logger.debug("Adding main donut to EBC six month view slide (13th slide)")
    left=Mm(43); top=Mm(48)
    slide_shapes.add_picture(file_donut_PA_ind_vols,left,top,height=Mm(115),width=Mm(94))

    #Update, by centre, the top interests, and the top industries with their top interests
    logger.debug("Setting the per-centre top 5 interests, together with per-centre top 3 industries and their interests")
    left_pos=[Mm(180),Mm(210),Mm(240),Mm(270),Mm(299)]   #distances to centre of icon for each row
    #First, the top interests for PA
    for col,p in enumerate(kwd_counts_6m_ctr["PA"].most_common(5)):
        # p is (keyword: count) for each of the keywords, so p[0] is the keyword itself
        replace_text_in_shape(slide_shapes,"PA-interest-{}".format(col),p[0],"13th slide")
        add_icon(slide_shapes,p[0],left=left_pos[col],top=Mm(59),small=True)
    #Next, the top industries with their interests for this centre - industry_counts_6m[c] is already ordered highest->lowest count
    n=0  #count how many displayed - need to do this separately from the loop count, as we ignore "Other" as an industry group
    for ind in industry_counts_6m["PA"].index:
        if (ind!="Other"):
            logger.debug("For Palo Alto, industry <%i> is <%s>" %(n,ind))
            #Write the list of top interests for this industry
            idx = find_text_in_shapes(slide_shapes,"PA-industry-{}".format(n))
            if (idx>=0):
                write_top_keywords(slide_shapes[idx].text_frame,
                                ind,
                                kwd_counts_6m_for_top_inds[("PA",ind)],
                                commented_rows_for_6m[("PA",ind)]
                                )
            else:
                logger.error("Could not find <PA-industry-{}> placeholder on 13th slide - idx={}".format(n,idx))
            n+=1   #increment the number of industries written out for this centre
        if (n>=3): break   #exit the loop after writing in 3 industries+interests

################################################
## 14th slide: 3 months of wordclouds based on keywords in objectives
################################################
logger.info(">>>> 14th slide: wordcloud for 3 months, based on Objectives")
df_for_month = dataframe_for_month(all_df, year=yyyy, month=mm)
kwd_count_for_month = kwds_in_objectives(df_for_month,vocab)
useful_rows_in_m = count_rows_with_comments(df_for_month)
logger.debug("Top keyword/counts for month %i : %r" % (mm,kwd_count_for_month.most_common(5)) )
file_wordcloud_for_month(kwd_count_for_month, useful_rows_in_m,
                         year=yyyy,month=mm)

df_for_month_minus_1 = dataframe_for_month(all_df, year=year_for_mm_minus_1, month=mm_minus_1)
kwd_count_for_m_minus_1 = kwds_in_objectives(df_for_month_minus_1,vocab)
logger.debug("Top keyword/counts for month %i : %r" % (mm_minus_1,kwd_count_for_m_minus_1.most_common(5)) )

df_for_month_minus_2 = dataframe_for_month(all_df, year=year_for_mm_minus_2, month=mm_minus_2)
kwd_count_for_m_minus_2 = kwds_in_objectives(df_for_month_minus_2,vocab)
logger.debug("Top keyword/counts for month %i : %r" % (mm_minus_2,kwd_count_for_m_minus_2.most_common(5)) )

useful_rows_in_m_2 = count_rows_with_comments(df_for_month_minus_2)
useful_rows_in_m_1 = count_rows_with_comments(df_for_month_minus_1)
logger.debug("Number of useful rows in months -2,-1,0 are %i, %i, %i" % (useful_rows_in_m_2,useful_rows_in_m_1,useful_rows_in_m))

file_wordcloud_for_month(kwd_count_for_m_minus_1, useful_rows_in_m_1, year=year_for_mm_minus_1,month=mm_minus_1)
file_wordcloud_for_month(kwd_count_for_m_minus_2, useful_rows_in_m_2, year=year_for_mm_minus_2,month=mm_minus_2)

## Modifying 3-month wordcloud slide by adding pic for last 3 months' wordcloud
logger.debug("Adding three wordclouds in 14th slide")
s = prs.slides[13]
slide_shapes=s.shapes
left=Mm(93)
top=Mm(36)
h=Mm(38)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_2]+".png",
                         left,top,height=h)
top=Mm(83)
slide_shapes.add_picture(tmpdir+"wordcloud-"+calendar.month_name[mm_minus_1]+".png",
                         left,top,height=h)
top=Mm(128)
slide_shapes.add_picture(tmpdir+"wordcloud-"+this_month+".png",
                         left,top,height=h)
#Find where the placeholders are for the keywords whose frequency we are graphing and update them
replace_text_in_shape(slide_shapes,find="Month-2",use=calendar.month_name[mm_minus_2],slidename="14th slide")
replace_text_in_shape(slide_shapes,find="Month-1",use=calendar.month_name[mm_minus_1],slidename="14th slide")
replace_text_in_shape(slide_shapes,find="Month-0",use=calendar.month_name[mm],slidename="14th slide")

notes_for_slide = s.notes_slide
notes_tf = notes_for_slide.notes_text_frame
notes_tf.text = ("%i rows with objectives in %s: Top keywords \n%r\n" %\
                    (useful_rows_in_m,this_month,kwd_count_for_month.most_common(10)) )
notes_tf.text += ("%i rows with objectives in %s: Top keywords \n%r\n" %\
                    (useful_rows_in_m_1,calendar.month_name[mm_minus_1],kwd_count_for_m_minus_1.most_common(10)) )
notes_tf.text += ("%i rows with objectives in %s: Top keywords \n%r\n" %\
                    (useful_rows_in_m_2,calendar.month_name[mm_minus_2],kwd_count_for_m_minus_2.most_common(10)) )

################################################
## Close the source presentations
################################################
logger.info("Saving Powerpoint file for "+this_month)
prs.save('GCA_Customer_Insights_'+this_month+'-'+str(yyyy)+'.pptx')
## Close any open figures
plt.close("all")
logger.info("...and we're done!")
for h in list(logger.handlers): logger.removeHandler(h)   # may be several here if we've crashed sometimes
